"""
OCT AI 프로젝트 발표 자료 자동 생성 스크립트 v3.
v2(20슬라이드) 기반으로 슬라이드 2개 추가 + 결과 슬라이드 업데이트.
결과: results/presentation_v3.pptx

추가/변경된 슬라이드:
  16. (NEW) 10~12단계: lr 조정 + 손실 함수 개선 탐색
  17. (NEW) 13단계: AROI N2N 사전학습 -> D1 fine-tuning
  18. 섹션4 구분 (업데이트)
  19. 종합 성능 비교 차트 (16개 방법으로 업데이트)
  20. 종합 성능 비교 테이블 (16개 방법으로 업데이트)
  21. 이미지 시각 비교 (유지)
  22. 핵심 발견 및 향후 계획 (업데이트)
"""

import io
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np

_korean_fonts = ["NanumGothic", "Malgun Gothic", "AppleGothic", "DejaVu Sans"]
for _fn in _korean_fonts:
    if any(_fn.lower() in f.name.lower() for f in fm.fontManager.ttflist):
        matplotlib.rcParams["font.family"] = _fn
        break

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT           = Path(__file__).parent.parent
RESULTS_DIR    = ROOT / "results"
IMG_BASELINE   = RESULTS_DIR / "01_baseline"   / "images"
IMG_SYNTH      = RESULTS_DIR / "02_synthetic_noise"
IMG_SUB2FULL   = RESULTS_DIR / "03_sub2full"   / "images"
IMG_SR_TEST    = RESULTS_DIR / "04_sr_test"
IMG_DNCNN      = RESULTS_DIR / "07_dncnn"      / "images" / "fold_1"
IMG_NAFNET     = RESULTS_DIR / "08_nafnet"     / "images" / "fold_1"
IMG_AUG        = RESULTS_DIR / "09_nafnet_aug" / "images" / "fold_1"
IMG_AROI_FT    = RESULTS_DIR / "13_aroi_n2n"   / "finetune" / "images" / "fold_1"
PRETRAIN_LOG   = RESULTS_DIR / "13_aroi_n2n"   / "pretrain" / "training_log.csv"

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
CANVAS   = RGBColor(0xfa, 0xf9, 0xf5)
CORAL    = RGBColor(0xcc, 0x78, 0x5c)
CORAL_D  = RGBColor(0xa9, 0x58, 0x3e)
DARK     = RGBColor(0x18, 0x17, 0x15)
DARK_ELV = RGBColor(0x25, 0x23, 0x20)
CARD     = RGBColor(0xef, 0xe9, 0xde)
INK      = RGBColor(0x14, 0x14, 0x13)
BODY     = RGBColor(0x3d, 0x3d, 0x3a)
MUTED    = RGBColor(0x6c, 0x6a, 0x64)
ON_DARK  = RGBColor(0xfa, 0xf9, 0xf5)
ON_DARK_S= RGBColor(0xa0, 0x9d, 0x96)
AMBER    = RGBColor(0xe8, 0xa5, 0x5a)
TEAL     = RGBColor(0x5d, 0xb8, 0xa6)
SUCCESS  = RGBColor(0x5d, 0xb8, 0x72)

SW = Inches(13.333)
SH = Inches(7.5)
SIZE_OFFSET = 3


# ---------------------------------------------------------------------------
# Utility (v2와 동일)
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill


def add_text(slide, text: str, x, y, w, h,
             size: int = 18, bold: bool = False, color: RGBColor = INK,
             align=PP_ALIGN.LEFT, font: str = "나눔고딕") -> None:
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size + SIZE_OFFSET)
    run.font.bold  = bold
    run.font.color.rgb = color


def add_img(slide, path: Path, x, y, w=None, h=None) -> None:
    if not path.exists():
        print(f"  [WARN] 이미지 없음: {path}")
        return
    if w and h:
        slide.shapes.add_picture(str(path), x, y, w, h)
    elif w:
        slide.shapes.add_picture(str(path), x, y, width=w)
    elif h:
        slide.shapes.add_picture(str(path), x, y, height=h)
    else:
        slide.shapes.add_picture(str(path), x, y)


def add_divider(slide, x, y, w, color: RGBColor = CORAL) -> None:
    add_rect(slide, x, y, w, Pt(2), color)


def buf_to_slide_img(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=120)
    buf.seek(0)
    plt.close(fig)
    return buf


def make_table_img(col_labels, rows, figsize, highlight_rows=None,
                   highlight_col=None, fontsize=10, scale_y=2.2):
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor("#faf9f5")
    ax.axis("off")
    table = ax.table(cellText=rows, colLabels=col_labels,
                     cellLoc="center", loc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(fontsize)
    table.scale(1, scale_y)
    table.auto_set_column_width(col=list(range(len(col_labels))))
    hl = set(highlight_rows) if highlight_rows else set()
    for (r, c), cell in table.get_celld().items():
        if r == 0:
            cell.set_facecolor("#181715")
            cell.get_text().set_color("#faf9f5")
            cell.get_text().set_fontweight("bold")
        elif r in hl:
            cell.set_facecolor("#e8ddd5")
            cell.get_text().set_fontweight("bold")
        elif r % 2 == 1:
            cell.set_facecolor("#efe9de")
        else:
            cell.set_facecolor("#faf9f5")
        if highlight_col is not None and c == highlight_col and r > 0:
            cell.get_text().set_color("#cc785c")
        cell.set_edgecolor("#e6dfd8")
    return buf_to_slide_img(fig)


def header_bar(slide, title: str, size: int = 28) -> None:
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.35), DARK)
    add_text(slide, title,
             Inches(0.6), Inches(0.28), Inches(12.0), Inches(0.8),
             size=size, bold=False, color=ON_DARK, font="나눔고딕")


# ---------------------------------------------------------------------------
# v2 슬라이드 함수들 (변경 없음)
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = blank_slide(prs)
    set_bg(s, DARK)
    add_rect(s, Inches(0.55), Inches(1.4), Inches(0.07), Inches(4.2), CORAL)
    add_text(s, "AI 기반 OCT 스페클 노이즈 제거",
             Inches(0.82), Inches(1.5), Inches(9.5), Inches(1.5),
             size=38, bold=False, color=ON_DARK, font="나눔고딕")
    add_text(s, "망막 선행 연구를 기반으로 한 치주 도메인 전이 전략",
             Inches(0.82), Inches(3.15), Inches(9.0), Inches(0.75),
             size=20, bold=False, color=ON_DARK_S, font="나눔고딕")
    add_divider(s, Inches(0.82), Inches(4.1), Inches(4.2), CORAL)
    add_text(s, "OCT Speckle Denoising — Retinal-to-Periodontal Transfer",
             Inches(0.82), Inches(4.3), Inches(8.5), Inches(0.6),
             size=13, color=ON_DARK_S, font="나눔고딕")
    add_rect(s, Inches(9.8), Inches(1.4), Inches(3.1), Inches(5.7), DARK_ELV)
    add_text(s, "발표 구성", Inches(10.0), Inches(1.6), Inches(2.7), Inches(0.4),
             size=12, color=CORAL, bold=True, font="나눔고딕")
    sections = [
        ("01  OCT 기초 및 평가 지표", CORAL),
        ("02  데이터", AMBER),
        ("03  방법론 (13단계)", TEAL),
        ("04  결과 비교", SUCCESS),
    ]
    for i, (txt, col) in enumerate(sections):
        add_rect(s, Inches(10.0), Inches(2.1 + i * 1.05), Inches(0.1), Inches(0.75), col)
        add_text(s, txt, Inches(10.22), Inches(2.12 + i * 1.05),
                 Inches(2.55), Inches(0.65), size=13, color=ON_DARK, font="나눔고딕")


def slide_section(prs, num, title, bullets, accent):
    s = blank_slide(prs)
    set_bg(s, DARK)
    add_rect(s, Inches(0.5), Inches(0), Inches(0.08), SH, accent)
    add_text(s, num, Inches(0.85), Inches(0.6), Inches(3.5), Inches(2.8),
             size=80, bold=False, color=accent, font="나눔고딕")
    add_text(s, title, Inches(0.85), Inches(3.55), Inches(9.5), Inches(1.4),
             size=34, bold=False, color=ON_DARK, font="나눔고딕")
    add_divider(s, Inches(0.85), Inches(5.05), Inches(9.2), accent)
    for i, b in enumerate(bullets):
        add_text(s, b, Inches(0.85), Inches(5.25 + i * 0.55),
                 Inches(9.5), Inches(0.48), size=14, color=ON_DARK_S, font="나눔고딕")


def slide_oct_basics(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "OCT 원리 및 스페클 노이즈")
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.6), Inches(5.6), CARD)
    add_text(s, "OCT 작동 원리", Inches(0.6), Inches(1.7), Inches(5.2), Inches(0.45),
             size=15, bold=True, color=CORAL, font="나눔고딕")
    add_text(s,
             "근적외선 마이켈슨 간섭계로 생체 조직 깊이 방향\n단면(B-scan)을 마이크로미터 단위로 재구성\n\n"
             "A-scan: 깊이 방향 1D 반사 신호\nB-scan: A-scan 연속 배열 => 2D 단면 이미지\n\n"
             "스페클 노이즈 원인:\n"
             "  레이저 가간섭성 => 후방 산란 위상 랜덤 간섭\n"
             "  촬영마다 독립적으로 다른 무작위 패턴 발생\n\n"
             "노이즈 모델 (곱셈성 + 가산성):\n",
             Inches(0.6), Inches(2.28), Inches(5.2), Inches(3.1),
             size=12, color=BODY, font="나눔고딕")
    add_rect(s, Inches(0.6), Inches(4.85), Inches(5.2), Inches(0.75), DARK_ELV)
    add_text(s, "I = S * Ns + Na\nNs ~ Gamma(L, 1/L)    Na ~ N(0, sigma^2)",
             Inches(0.7), Inches(4.92), Inches(5.0), Inches(0.65),
             size=11, color=ON_DARK, font="Courier New")
    add_text(s,
             "I: 합성 noisy 이미지    S: 원본 clean\n"
             "Ns: 곱셈성 스페클 (mean=1, var=1/L)    Na: 가산성 가우시안\n"
             "SBSDI D1 캘리브레이션:  L = 5.27,  sigma = 0.010",
             Inches(0.6), Inches(5.68), Inches(5.2), Inches(1.05),
             size=10, color=BODY, font="나눔고딕")
    add_text(s, "실제 OCT B-scan 예시 (SBSDI D1 #01)",
             Inches(6.3), Inches(1.58), Inches(6.7), Inches(0.42),
             size=13, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_BASELINE / "sample_01.png", Inches(6.2), Inches(2.1), w=Inches(6.9))
    add_text(s, "좌->우: noisy 단일 프레임 / clean 다중프레임 평균 / NLM / BM3D / SRAD",
             Inches(6.2), Inches(6.72), Inches(6.9), Inches(0.42),
             size=10, color=MUTED, font="나눔고딕")


def slide_metrics(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "평가 지표 — PSNR / SSIM / CNR")
    items = [
        ("PSNR", "Peak Signal-to-Noise Ratio\n최대 신호 대 잡음비", CORAL,
         "PSNR = 10 * log10( MAX^2 / MSE )",
         "MSE: 복원 vs clean 픽셀 제곱 오차 평균\nMAX: 픽셀 최댓값 (정규화 기준 1.0)\n\n"
         "오차 작을수록 PSNR 높음 => 클수록 좋음\n단위: dB, 실용 범위 20~40 dB\n\n"
         "한계: 픽셀 수치 오차만 반영 (구조 무시)"),
        ("SSIM", "Structural Similarity Index\n구조적 유사도", AMBER,
         "SSIM = f(밝기, 대비, 구조 공분산)",
         "밝기(mu) + 대비(sigma^2) + 구조(sigma_xy)\n동시 비교 => 인간 시각에 가까운 평가\n\n"
         "범위: 0~1, 높을수록 좋음\nPSNR보다 임상 판독 품질에 가까운 지표\n\n"
         "한계: 지역 윈도(11x11) 기반"),
        ("CNR", "Contrast-to-Noise Ratio\n대비 대 잡음비", TEAL,
         "CNR = |mu_sig - mu_bg| / sigma_bg",
         "OCT 조직 레이어(신호)와 배경(잡음) 간\n대비를 배경 잡음 수준으로 나눈 값\n\n"
         "높을수록 레이어 경계 선명, 구조 식별 용이\n\n"
         "임상적 진단 품질에 직접 대응하는\nOCT 특화 지표 — 딥러닝에서 가장 달성 어려움"),
    ]
    for i, (abbr, full, col, formula, desc) in enumerate(items):
        cx = Inches(0.4 + i * 4.28)
        add_rect(s, cx, Inches(1.5), Inches(4.1), Inches(5.6), CARD)
        add_rect(s, cx, Inches(1.5), Inches(4.1), Inches(0.1), col)
        add_text(s, abbr, cx + Inches(0.2), Inches(1.72), Inches(3.7), Inches(0.62),
                 size=22, bold=True, color=INK, font="나눔고딕")
        add_text(s, full, cx + Inches(0.2), Inches(2.38), Inches(3.7), Inches(0.55),
                 size=11, color=MUTED, font="나눔고딕")
        add_divider(s, cx + Inches(0.2), Inches(3.0), Inches(3.5), col)
        add_rect(s, cx + Inches(0.2), Inches(3.1), Inches(3.7), Inches(0.62), DARK_ELV)
        add_text(s, formula, cx + Inches(0.3), Inches(3.18), Inches(3.5), Inches(0.5),
                 size=10, color=ON_DARK, font="Courier New")
        add_text(s, desc, cx + Inches(0.2), Inches(3.82), Inches(3.7), Inches(3.1),
                 size=11, color=BODY, font="나눔고딕")


def slide_data_overview(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "보유 데이터셋 현황")
    add_text(s, "SBSDI 패키지 내부 구조 (Fang et al., TMI 2013)",
             Inches(0.4), Inches(1.52), Inches(8.0), Inches(0.42),
             size=13, bold=True, color=CORAL, font="나눔고딕")
    sbsdi_cols = ["서브디렉토리", "명칭", "폴더수", "파일 구성", "GT"]
    sbsdi_rows = [
        ["synthetic experiments/", "D1", "18", "test.tif, average.tif, 1~4.tif", "있음 (average.tif)"],
        ["real experiments (Humans)/", "D2", "39", "test.tif, 1~4.tif", "없음"],
        ["real experiments (Mouse)/", "D3", "3 (해상도별)", "test.tif, 1~4.tif", "없음"],
    ]
    buf = make_table_img(sbsdi_cols, sbsdi_rows, figsize=(12.5, 2.2),
                         highlight_rows={1}, fontsize=9, scale_y=2.0)
    s.shapes.add_picture(buf, Inches(0.4), Inches(1.98), width=Inches(12.6))
    add_rect(s, Inches(0.4), Inches(3.48), Inches(5.8), Inches(0.75), DARK_ELV)
    add_text(s, "average.tif 생성 방법: 동일 위치 N~40회 반복 스캔 => 픽셀별 평균 => clean GT",
             Inches(0.6), Inches(3.58), Inches(5.4), Inches(0.58),
             size=11, bold=True, color=ON_DARK, font="나눔고딕")
    add_rect(s, Inches(6.5), Inches(3.48), Inches(6.45), Inches(0.75), DARK_ELV)
    add_text(s, "1~4.tif는 인접 공간 슬라이스 (다른 위치 B-scan) — 같은 위치 반복 스캔이 아님",
             Inches(6.65), Inches(3.58), Inches(6.15), Inches(0.58),
             size=11, bold=True, color=AMBER, font="나눔고딕")
    add_text(s, "추가 보유 데이터셋",
             Inches(0.4), Inches(4.38), Inches(8.0), Inches(0.42),
             size=13, bold=True, color=CORAL, font="나눔고딕")
    other_cols = ["데이터셋", "규모", "태스크", "GT", "활용"]
    other_rows = [
        ["AROI\n(Heidelberg)", "3,072장\n24명 / 128 B-scan/명", "레이어\n세그멘테이션", "없음\n-> 합성/N2N",
         "합성 학습 1,136쌍\n13단계 N2N 사전학습 3,048쌍"],
        ["Kermany OCT2017\n(CC BY 4.0)", "84,484장\n4클래스 분류", "질환 분류", "없음\n-> 합성",
         "합성 학습 5,000쌍\n랜덤 샘플"],
    ]
    buf2 = make_table_img(other_cols, other_rows, figsize=(12.5, 2.4),
                          fontsize=9.5, scale_y=2.4)
    s.shapes.add_picture(buf2, Inches(0.4), Inches(4.85), width=Inches(9.3))
    add_text(s, "AROI 합성 노이즈 샘플",
             Inches(9.95), Inches(4.42), Inches(3.2), Inches(0.38),
             size=11, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_SYNTH / "AROI" / "samples" / "sample_00000.png",
            Inches(9.95), Inches(4.85), w=Inches(3.1))


def slide_data_survey(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "외부 데이터 조사 및 clean GT 생성 가능성 분석")
    add_text(s, "Duke Farsiu 연구실 공개 데이터 탐색 (people.duke.edu/~sf59, 2026-06-01)",
             Inches(0.4), Inches(1.52), Inches(12.5), Inches(0.42),
             size=13, bold=True, color=CORAL, font="나눔고딕")
    duke_cols = ["데이터셋", "규모", "GT", "비고"]
    duke_rows = [
        ["SBSDI D1 (이미 보유)", "18쌍", "있음", "유일한 GT 보유 데이터셋"],
        ["AMD SD-OCT (RPEDC 2013)", "384명 / 38,400 B-scan", "없음", "AMD vs 정상 분류용 / 자가지도학습에만 활용"],
        ["Srinivasan BOE 2014", "45명 (정상/AMD/DME 각 15명)", "없음", "소규모, 분류용"],
        ["Chiu DME BOE 2014", "다수", "없음", "낭종 세그멘테이션용"],
    ]
    buf = make_table_img(duke_cols, duke_rows, figsize=(12.5, 2.2),
                         highlight_rows={1}, fontsize=9.5, scale_y=2.0)
    s.shapes.add_picture(buf, Inches(0.4), Inches(1.98), width=Inches(12.6))
    add_text(s, "보유 데이터 clean GT 추가 생성 가능 여부 (diff/frame std 비율 분석)",
             Inches(0.4), Inches(3.62), Inches(12.5), Inches(0.42),
             size=13, bold=True, color=AMBER, font="나눔고딕")
    add_rect(s, Inches(0.4), Inches(4.12), Inches(7.0), Inches(0.62), DARK_ELV)
    add_text(s, "판별 기준: diff(A-B) std / A std — 같은 위치 반복 스캔이면 ~0.77 (구조 상쇄), 다른 위치이면 >=1.0",
             Inches(0.55), Inches(4.22), Inches(6.7), Inches(0.48),
             size=10, color=ON_DARK, font="나눔고딕")
    gt_cols = ["데이터셋", "비율", "판단", "GT 생성"]
    gt_rows = [
        ["SBSDI D1 (test vs avg)", "0.771", "동일위치 반복스캔", "GT 이미 존재"],
        ["SBSDI D2 (test vs 1~4)", "1.07", "인접공간 슬라이스", "불가"],
        ["AROI (순차 B-scan)", "0.86~1.31", "3D 순차 슬라이스", "불가 -> N2N 활용"],
        ["Kermany OCT2017", "—", "세션/해상도 혼재", "불가"],
    ]
    buf2 = make_table_img(gt_cols, gt_rows, figsize=(7.2, 2.8),
                          highlight_rows={1}, fontsize=9.5, scale_y=2.0)
    s.shapes.add_picture(buf2, Inches(0.4), Inches(4.82), width=Inches(7.1))
    add_rect(s, Inches(7.8), Inches(4.12), Inches(5.1), Inches(3.1), DARK_ELV)
    add_text(s, "결론", Inches(8.0), Inches(4.28), Inches(4.7), Inches(0.38),
             size=12, bold=True, color=CORAL, font="나눔고딕")
    add_text(s,
             "현재 보유 데이터 중 clean GT를\n추가 생성 가능한 데이터셋 없음.\n\n"
             "외부 접근 실패 데이터셋:\n"
             "  PKU37 (37쌍) — 알리바바 클라우드\n"
             "  RETOUCH (112 볼륨) — Grand Challenge\n"
             "  ODTiD (242장) — OPENICPSR\n\n"
             "근본적 해결 방향:\n"
             "  AROI N2N 사전학습 (13단계)\n"
             "  치주 OCT 직접 촬영",
             Inches(8.0), Inches(4.72), Inches(4.7), Inches(2.4),
             size=11, color=ON_DARK_S, font="나눔고딕")


def slide_pipeline_overview(prs):
    s = blank_slide(prs)
    set_bg(s, DARK)
    add_text(s, "전체 실험 파이프라인 개요", Inches(0.6), Inches(0.22),
             Inches(11), Inches(0.78), size=28, bold=False, color=ON_DARK, font="나눔고딕")
    add_divider(s, Inches(0.6), Inches(1.05), Inches(12.1), CORAL)
    left_cols  = ["단계", "방법/모델", "핵심 내용", "대표 결과"]
    right_cols = ["단계", "방법/모델", "핵심 내용", "대표 결과"]
    left_rows = [
        ["1단계", "NLM/BM3D\n/SRAD", "전통 신호처리\n베이스라인", "SRAD\nPSNR 27.50"],
        ["2단계", "합성 노이즈\n생성", "Gamma+Gaussian\n6,136쌍 생성", "KS=0.119"],
        ["3단계-A", "N2N\n자가지도", "real 프레임쌍\n468쌍/GT 불필요", "SSIM\n0.681(+)"],
        ["3단계-B", "지도학습\nU-Net", "합성 6,136쌍\n=> 도메인 갭", "PSNR\n22.43(-)"],
        ["4단계", "Real-ESRGAN\nx2/x4", "Blind SR 사전학습\n노이즈+SR 동시", "PSNR\n27.57(+)"],
        ["5단계", "Pre-train\n+Fine-tune", "합성 사전학습\n+N2N FT", "PSNR\n27.46"],
        ["6단계", "6-fold CV\nU-Net", "real clean GT\nbatch=64 ES", "PSNR\n28.19(+)"],
    ]
    right_rows = [
        ["7단계", "DnCNN\n667K scratch", "6-fold CV\n아키텍처 비교", "PSNR\n28.17"],
        ["8단계", "NAFNet-32\n17M scratch", "SimpleGate\n+SCA 어텐션", "CNR\n1.183"],
        ["9단계", "NAFNet\n+Aug K=4", "노이즈 재실현\n5배 증강", "PSNR\n28.35(+)"],
        ["10단계", "NAFNet\nlr=1e-4", "lr 조정 검증\n=> 오히려 하락", "PSNR\n27.95(-)"],
        ["11단계", "NAFNet\n+EdgeLoss", "Sobel 경계 손실\nλ=0.5 안정화", "CNR\n1.186"],
        ["12단계", "NAFNet\n+Edge+Freq", "FFT 주파수 손실\n추가 => 역효과", "CNR\n1.182(-)"],
        ["13단계", "AROI N2N\n사전학습+FT", "3,048쌍 사전학습\n=> D1 fine-tuning", "SSIM\n0.6838(+)"],
    ]
    buf_l = make_table_img(left_cols, left_rows, figsize=(6.2, 4.5), fontsize=9, scale_y=2.0)
    buf_r = make_table_img(right_cols, right_rows, figsize=(6.2, 4.5), fontsize=9, scale_y=2.0)
    s.shapes.add_picture(buf_l, Inches(0.4), Inches(1.2), width=Inches(6.2))
    s.shapes.add_picture(buf_r, Inches(6.8), Inches(1.2), width=Inches(6.2))
    add_rect(s, Inches(0.4), Inches(6.32), Inches(12.5), Inches(0.75), DARK_ELV)
    add_text(s,
             "(+): SRAD 해당 지표 초과  |  (-): 미달  |  6~13단계: 6-fold CV 체계적 비교  |  13단계 SSIM 전체 최고",
             Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.55),
             size=11, color=ON_DARK_S, align=PP_ALIGN.CENTER, font="나눔고딕")


def slide_traditional(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "1단계: 전통적 방법 베이스라인 (NLM / BM3D / SRAD)")
    methods = [
        ("NLM", "Non-Local Means", CORAL,
         "이미지 전체에서 유사 패치를 찾아\n가중 평균으로 노이즈 제거\n속도: ~0.51 s/장"),
        ("BM3D", "Block-Matching 3D", AMBER,
         "유사 블록을 3D 변환 도메인으로\n묶어 임계화 후 재구성\n속도: ~3.20 s/장"),
        ("SRAD", "Speckle Reducing\nAnisotropic Diffusion", TEAL,
         "곱셈성 스페클 모델 반영 PDE 기반\n이방성 확산, 경계 보존 최우수\n속도: ~5.68 s/장"),
    ]
    for i, (abbr, full, col, desc) in enumerate(methods):
        cx = Inches(0.4 + i * 3.08)
        add_rect(s, cx, Inches(1.5), Inches(2.88), Inches(2.4), CARD)
        add_rect(s, cx, Inches(1.5), Inches(2.88), Inches(0.08), col)
        add_text(s, abbr, cx + Inches(0.15), Inches(1.65), Inches(2.6), Inches(0.46),
                 size=16, bold=True, color=INK, font="나눔고딕")
        add_text(s, full, cx + Inches(0.15), Inches(2.15), Inches(2.6), Inches(0.45),
                 size=10, color=MUTED, font="나눔고딕")
        add_text(s, desc, cx + Inches(0.15), Inches(2.65), Inches(2.6), Inches(1.15),
                 size=11, color=BODY, font="나눔고딕")
    col_labels = ["방법", "PSNR (dB)", "SSIM", "CNR", "처리 속도"]
    rows = [
        ["NLM",  "26.12+-2.01", "0.492+-0.050", "1.130+-0.119", "0.51 s/장"],
        ["BM3D", "27.00+-2.51", "0.599+-0.066", "1.134+-0.119", "3.20 s/장"],
        ["SRAD", "27.50+-1.98", "0.652+-0.023", "1.220+-0.121", "5.68 s/장"],
    ]
    buf = make_table_img(col_labels, rows, figsize=(5.0, 1.9),
                         highlight_rows={3}, fontsize=9.5, scale_y=2.2)
    s.shapes.add_picture(buf, Inches(0.4), Inches(4.08), width=Inches(5.2))
    add_rect(s, Inches(0.4), Inches(6.48), Inches(5.2), Inches(0.75), CORAL)
    add_text(s, "SRAD: PSNR / SSIM / CNR 세 지표 모두 1위 => 딥러닝 목표 기준값",
             Inches(0.6), Inches(6.6), Inches(4.8), Inches(0.55),
             size=12, bold=True, color=ON_DARK, font="나눔고딕")
    add_text(s, "시각적 비교 (sample #01: Noisy / Clean / NLM / BM3D / SRAD)",
             Inches(5.75), Inches(3.85), Inches(7.3), Inches(0.45),
             size=11, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_BASELINE / "sample_01.png", Inches(5.75), Inches(4.35), w=Inches(7.3))


def slide_synthetic_noise(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "2단계: 합성 스페클 노이즈 생성 파이프라인")
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.45), Inches(5.6), CARD)
    add_text(s, "물리 기반 노이즈 모델", Inches(0.6), Inches(1.68),
             Inches(5.0), Inches(0.42), size=14, bold=True, color=CORAL, font="나눔고딕")
    add_rect(s, Inches(0.6), Inches(2.18), Inches(5.0), Inches(0.78), DARK_ELV)
    add_text(s, "I = S * Ns + Na\nNs ~ Gamma(L, 1/L)    Na ~ N(0, sigma^2)",
             Inches(0.7), Inches(2.25), Inches(4.8), Inches(0.68),
             size=11, color=ON_DARK, font="Courier New")
    add_text(s,
             "I: 합성 noisy 이미지    S: 원본 clean 이미지\n"
             "Ns: 곱셈성 스페클 (mean=1, var=1/L)\n"
             "Na: 가산성 가우시안 노이즈\n"
             "L: looks 수 — 높을수록 노이즈 약함\n\n"
             "주의: AROI/Kermany는 real OCT 이미지로\n"
             "이미 스페클을 포함. 원본을 의사(pseudo)\n"
             "clean으로 삼아 합성 스페클 추가.",
             Inches(0.6), Inches(3.08), Inches(5.0), Inches(1.85),
             size=12, color=BODY, font="나눔고딕")
    add_divider(s, Inches(0.6), Inches(5.02), Inches(4.9), AMBER)
    add_text(s, "파라미터 캘리브레이션 (SBSDI D1 18쌍 기반)",
             Inches(0.6), Inches(5.12), Inches(4.9), Inches(0.42),
             size=12, bold=True, color=INK, font="나눔고딕")
    add_text(s,
             "  L = 5.266 +- 0.743  /  sigma = 0.010\n"
             "  KS 통계량 = 0.119",
             Inches(0.6), Inches(5.6), Inches(5.0), Inches(0.75),
             size=12, color=BODY, font="나눔고딕")
    add_text(s, "합성 노이즈 샘플 (AROI)", Inches(6.1), Inches(1.58),
             Inches(7.0), Inches(0.42), size=13, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_SYNTH / "AROI" / "samples" / "sample_00000.png",
            Inches(6.1), Inches(2.05), w=Inches(7.0))
    add_rect(s, Inches(6.1), Inches(5.65), Inches(7.0), Inches(1.45), DARK_ELV)
    add_text(s, "생성된 학습 데이터", Inches(6.3), Inches(5.78),
             Inches(6.6), Inches(0.38), size=11, bold=True, color=CORAL, font="나눔고딕")
    add_text(s,
             "AROI         1,136쌍   512x1024px   90도 회전 보정\n"
             "Kermany  5,000쌍   496x512px    랜덤 샘플링 (seed=42)\n"
             "합계          6,136쌍   재현 가능 파이프라인",
             Inches(6.3), Inches(6.22), Inches(6.6), Inches(0.85),
             size=11, color=ON_DARK, font="나눔고딕")


def slide_n2n_supervised(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "3단계: 자가지도 학습 (N2N) + 지도학습 (도메인 갭 분석)")
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.6), CARD)
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(0.1), TEAL)
    add_text(s, "3단계-A: Noise2Noise 자가지도 학습",
             Inches(0.6), Inches(1.7), Inches(5.5), Inches(0.45),
             size=14, bold=True, color=TEAL, font="나눔고딕")
    add_text(s,
             "clean GT 없이 OCT 이미지만으로 학습.\n\n"
             "SBSDI D2 39세트 × 12 순서쌍 = 468쌍\n"
             "(frame_i, frame_j) => noisy-noisy 학습 쌍\n\n"
             "U-Net (1.95M params)  /  L1 Loss\n"
             "lr=1e-4 CosineAnnealing  /  500 epoch\n"
             "학습 시간: ~221분 (RTX A4000)\n\n"
             "PSNR  26.84 +- 2.18  dB  (SRAD: 27.50)\n"
             "SSIM   0.681 +- 0.033    (SRAD:  0.652)  [+]\n"
             "CNR    1.148 +- 0.133    (SRAD:  1.220)",
             Inches(0.6), Inches(2.25), Inches(5.5), Inches(3.0),
             size=12, color=BODY, font="나눔고딕")
    add_img(s, IMG_SUB2FULL / "sample_01.png",
            Inches(0.45), Inches(5.22), w=Inches(5.8))
    add_rect(s, Inches(6.65), Inches(1.5), Inches(6.25), Inches(5.6), CARD)
    add_rect(s, Inches(6.65), Inches(1.5), Inches(6.25), Inches(0.1), CORAL)
    add_text(s, "3단계-B: 지도학습 (합성 데이터 — 도메인 갭)",
             Inches(6.85), Inches(1.7), Inches(5.85), Inches(0.45),
             size=13, bold=True, color=CORAL, font="나눔고딕")
    add_text(s,
             "학습: 6,136쌍 (합성 Gamma+Gaussian noisy)\n"
             "평가: SBSDI D1 18쌍 (실제 OCT 스페클)\n\n"
             "U-Net (동일 아키텍처)  /  L1 Loss\n"
             "lr=2e-4  /  100 epoch  /  ~305분\n"
             "Best loss: 0.021621 (epoch 10 이후 포화)\n\n"
             "PSNR  22.43  dB   (SRAD: 27.50)  [-5.07]\n"
             "SSIM   0.333       (SRAD:  0.652)  [-0.319]\n"
             "CNR    0.902       (SRAD:  1.220)  [-0.318]",
             Inches(6.85), Inches(2.25), Inches(5.85), Inches(3.0),
             size=12, color=BODY, font="나눔고딕")
    add_rect(s, Inches(6.85), Inches(5.38), Inches(5.85), Inches(1.38), DARK_ELV)
    add_text(s, "도메인 갭 원인 분석",
             Inches(7.0), Inches(5.5), Inches(5.5), Inches(0.38),
             size=11, bold=True, color=AMBER, font="나눔고딕")
    add_text(s,
             "학습 노이즈: 합성 Gamma+Gaussian (수학 모델)\n"
             "평가 노이즈: 실제 OCT 스페클 (복잡한 물리)\n"
             "=> 노이즈 모델 정확도가 지도학습 성패를 좌우",
             Inches(7.0), Inches(5.9), Inches(5.5), Inches(0.82),
             size=11, color=ON_DARK_S, font="나눔고딕")


def slide_sr_finetune(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "4~5단계: Real-ESRGAN SR 적용 + Pre-train Fine-tune")
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.6), CARD)
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(0.1), TEAL)
    add_text(s, "4단계: Real-ESRGAN Blind SR",
             Inches(0.6), Inches(1.7), Inches(5.5), Inches(0.45),
             size=14, bold=True, color=TEAL, font="나눔고딕")
    add_text(s,
             "Blind SR: 노이즈+블러+압축 등 복합 열화\n동시 복원하는 사전 학습 모델\n\n"
             "RRDBNet (23 잔차 밀집 블록)\nRealESRGAN_x2plus.pth / x4plus.pth (64MB)\nCUDA half precision, tile=400\n",
             Inches(0.6), Inches(2.28), Inches(5.5), Inches(1.75),
             size=12, color=BODY, font="나눔고딕")
    sr_cols = ["방법", "PSNR (dB)", "SSIM", "CNR", "속도"]
    sr_rows = [
        ["SRAD (기준)", "27.50+-1.98", "0.652+-0.023", "1.220", "5.68 s"],
        ["ESRGAN x2",  "27.30+-2.35", "0.674+-0.034", "1.121", "0.62 s"],
        ["ESRGAN x4",  "27.57+-2.40", "0.673+-0.034", "1.166", "1.42 s"],
    ]
    buf = make_table_img(sr_cols, sr_rows, figsize=(6.0, 1.8),
                         highlight_rows={3}, fontsize=9.5, scale_y=2.2)
    s.shapes.add_picture(buf, Inches(0.4), Inches(4.12), width=Inches(5.9))
    add_text(s, "x4: PSNR+SSIM 모두 SRAD 초과 / CNR 미달 / 처리 속도 4배",
             Inches(0.4), Inches(5.82), Inches(5.9), Inches(0.38),
             size=10, bold=True, color=TEAL, font="나눔고딕")
    add_img(s, IMG_SR_TEST / "all_comparison.png",
            Inches(0.4), Inches(6.25), w=Inches(5.9))
    add_rect(s, Inches(6.65), Inches(1.5), Inches(6.25), Inches(5.6), CARD)
    add_rect(s, Inches(6.65), Inches(1.5), Inches(6.25), Inches(0.1), AMBER)
    add_text(s, "5단계: Pre-train + Fine-tune (N2N)",
             Inches(6.85), Inches(1.7), Inches(5.85), Inches(0.45),
             size=14, bold=True, color=AMBER, font="나눔고딕")
    add_text(s,
             "합성 지도학습 가중치(3단계-B)로 초기화\n=> N2N real 쌍으로 fine-tuning\n\n"
             "Loss: L1 + 0.1*(1-SSIM)  /  lr=1e-5\n"
             "에포크: 200 / 배치: 16 / ~35분\n"
             "Best loss: 0.184206 (epoch 20 이후 포화)",
             Inches(6.85), Inches(2.28), Inches(5.85), Inches(1.88),
             size=12, color=BODY, font="나눔고딕")
    ft_cols = ["방법", "PSNR (dB)", "SSIM", "CNR"]
    ft_rows = [
        ["SRAD",           "27.50+-1.98", "0.652+-0.023", "1.220+-0.121"],
        ["N2N (scratch)",  "26.84+-2.18", "0.681+-0.033", "1.148+-0.133"],
        ["Pre-train + FT", "27.46+-2.55", "0.679+-0.037", "1.169+-0.120"],
    ]
    buf2 = make_table_img(ft_cols, ft_rows, figsize=(6.2, 1.8),
                          highlight_rows={3}, fontsize=9.5, scale_y=2.2)
    s.shapes.add_picture(buf2, Inches(6.65), Inches(4.25), width=Inches(6.2))
    add_rect(s, Inches(6.85), Inches(5.92), Inches(5.85), Inches(1.1), DARK_ELV)
    add_text(s,
             "전이학습 효과 확인: PSNR +0.62 dB (N2N 대비)\n"
             "SRAD까지 0.04 dB 차이로 사실상 동급\n"
             "N2N 노이즈 타겟 분산이 loss floor(~0.184) 형성",
             Inches(7.05), Inches(6.02), Inches(5.5), Inches(0.95),
             size=11, color=ON_DARK_S, font="나눔고딕")


def slide_kfold_archs(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "6~8단계: 6-fold CV 아키텍처 비교 (U-Net / DnCNN / NAFNet)")
    add_text(s, "동일 프레임워크: SBSDI D1 18쌍 6-fold CV / real clean GT / L1+0.1*(1-SSIM) / Early stopping",
             Inches(0.4), Inches(1.48), Inches(12.5), Inches(0.38),
             size=11, bold=True, color=CORAL, font="나눔고딕")
    arch_cols = ["단계", "아키텍처", "params", "초기화", "학습시간", "PSNR (dB)", "SSIM", "CNR"]
    arch_rows = [
        ["6단계-B", "U-Net\n(3단계 공유)", "1.95M", "합성\n사전학습", "47분", "28.19\n+-2.56", "0.6814", "1.169"],
        ["7단계", "DnCNN-B\n(depth=20)", "667K", "from\nscratch", "355분", "28.17\n+-2.47", "0.6732", "1.167"],
        ["8단계", "NAFNet-32", "17M", "from\nscratch", "123분", "28.11\n+-2.26", "0.6743", "1.183"],
    ]
    buf = make_table_img(arch_cols, arch_rows, figsize=(12.8, 2.5),
                         highlight_rows={1}, fontsize=9.5, scale_y=2.4)
    s.shapes.add_picture(buf, Inches(0.4), Inches(1.92), width=Inches(12.6))
    add_rect(s, Inches(0.4), Inches(3.75), Inches(12.6), Inches(0.72), DARK_ELV)
    add_text(s,
             "핵심 발견: 667K(DnCNN) = 1.95M(U-Net) = 17M(NAFNet) — 파라미터 수 25배 차이에도 PSNR 0.08 dB 이내",
             Inches(0.6), Inches(3.88), Inches(12.2), Inches(0.5),
             size=12, bold=True, color=ON_DARK, align=PP_ALIGN.CENTER, font="나눔고딕")
    add_text(s, "DnCNN fold_1 복원 결과 (sample_01)",
             Inches(0.4), Inches(4.58), Inches(5.9), Inches(0.38),
             size=11, bold=True, color=BODY, font="나눔고딕")
    add_img(s, IMG_DNCNN / "sample_01.png", Inches(0.4), Inches(5.0), w=Inches(5.9))
    add_text(s, "NAFNet fold_1 복원 결과 (sample_01)",
             Inches(6.65), Inches(4.58), Inches(6.2), Inches(0.38),
             size=11, bold=True, color=BODY, font="나눔고딕")
    add_img(s, IMG_NAFNET / "sample_01.png", Inches(6.65), Inches(5.0), w=Inches(6.2))
    add_rect(s, Inches(0.4), Inches(6.78), Inches(12.6), Inches(0.62), CARD)
    add_text(s,
             "NAFNet: CNR 1.183으로 AI 방법 중 최고 / 학습 시간 DnCNN 대비 1/3 / fold 간 분산 최소(std 2.26)",
             Inches(0.6), Inches(6.88), Inches(12.2), Inches(0.42),
             size=11, color=BODY, align=PP_ALIGN.CENTER, font="나눔고딕")


def slide_augmentation(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "9단계: 다중 노이즈 재실현 증강 + NAFNet 6-fold CV")
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.6), CARD)
    add_text(s, "증강 전략", Inches(0.6), Inches(1.68),
             Inches(5.5), Inches(0.42), size=14, bold=True, color=CORAL, font="나눔고딕")
    add_text(s,
             "동기: 8단계까지 3번 연속 데이터 병목 확인\n18쌍 크기가 아키텍처 차이를 완전히 압도\n\n"
             "전략: 18개 clean GT에 노이즈 모델(L=5.266)\n"
             "로 K=4개씩 합성 noisy 재실현 생성\n=> 학습 데이터 1배 => 5배 (15+60=75쌍/fold)\n\n"
             "평가: real noisy 3쌍/fold (8단계와 동등 조건)\n\n"
             "NAFNet-32 / batch=48 / steps_per_epoch=425\n"
             "Early stopping patience=30",
             Inches(0.6), Inches(2.22), Inches(5.5), Inches(3.35),
             size=12, color=BODY, font="나눔고딕")
    aug_cols = ["방법", "PSNR (dB)", "SSIM", "CNR"]
    aug_rows = [
        ["SRAD",           "27.50+-1.98", "0.652+-0.023", "1.220+-0.121"],
        ["NAFNet (8단계)", "28.11+-2.26", "0.6743+-0.029", "1.183+-0.120"],
        ["NAFNet+Aug K=4", "28.35+-2.56", "0.6832+-0.032", "1.168+-0.121"],
    ]
    buf = make_table_img(aug_cols, aug_rows, figsize=(5.9, 1.8),
                         highlight_rows={3}, fontsize=9.5, scale_y=2.2)
    s.shapes.add_picture(buf, Inches(0.4), Inches(5.65), width=Inches(5.9))
    add_text(s, "NAFNet+Aug fold_1 복원 결과 (sample_01)",
             Inches(6.65), Inches(1.55), Inches(6.2), Inches(0.42),
             size=11, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_AUG / "sample_01.png", Inches(6.65), Inches(2.0), w=Inches(6.2))
    add_rect(s, Inches(6.65), Inches(5.38), Inches(6.2), Inches(1.72), DARK_ELV)
    add_text(s, "결과 분석", Inches(6.85), Inches(5.5),
             Inches(5.8), Inches(0.38), size=12, bold=True, color=CORAL, font="나눔고딕")
    add_text(s,
             "PSNR +0.24 dB / SSIM +0.009: 5배 데이터의 효과\n"
             "CNR -0.015: 합성 노이즈 != 실제 스페클 (오차 범위 내)\n"
             "Early stop 31~34 에포크 (8단계 32~46과 유사)\n"
             "=> 데이터 병목 부분 완화, 외부 real OCT 확보 필요",
             Inches(6.85), Inches(5.93), Inches(5.8), Inches(1.1),
             size=11, color=ON_DARK_S, font="나눔고딕")


# ---------------------------------------------------------------------------
# 슬라이드 16 (NEW): 10~12단계 lr 조정 + 손실 함수 개선
# ---------------------------------------------------------------------------

def slide_lr_loss_exp(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "10~12단계: 하이퍼파라미터 & 손실 함수 개선 탐색")

    # 왼쪽: lr 조정 (10단계)
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.6), CARD)
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(0.1), AMBER)
    add_text(s, "10단계: lr=1e-4 조정 실험",
             Inches(0.6), Inches(1.72), Inches(5.5), Inches(0.42),
             size=14, bold=True, color=AMBER, font="나눔고딕")
    add_text(s,
             "가설: lr=1e-3이 너무 높아 epoch 1~9에서\n"
             "val PSNR 최고점 후 과적합 => lr 낮추면?\n\n"
             "Epoch 1 val PSNR 비교:\n"
             "  lr=1e-3: ~28.27 (이미 높음)\n"
             "  lr=1e-4: ~21~25 (낮은 시작)\n\n"
             "lr=1e-4는 patience=30 내에 lr=1e-3의\n"
             "초반 수렴점까지 도달하지 못함.",
             Inches(0.6), Inches(2.28), Inches(5.5), Inches(2.45),
             size=12, color=BODY, font="나눔고딕")
    lr_cols = ["lr", "PSNR", "SSIM", "CNR", "결과"]
    lr_rows = [
        ["1e-3 (8단계)", "28.11", "0.6743", "1.183", "기준"],
        ["1e-4 (10단계)", "27.95", "0.6666", "1.191", "PSNR -0.16 (하락)"],
    ]
    buf = make_table_img(lr_cols, lr_rows, figsize=(5.9, 1.5),
                         highlight_rows={2}, fontsize=9, scale_y=2.2)
    s.shapes.add_picture(buf, Inches(0.4), Inches(5.0), width=Inches(5.9))
    add_rect(s, Inches(0.4), Inches(6.38), Inches(5.9), Inches(0.72), DARK_ELV)
    add_text(s, "결론: 데이터 18쌍 환경에서 빠른 초기 수렴이 유리. lr 조정으로 데이터 병목 해소 불가.",
             Inches(0.55), Inches(6.48), Inches(5.6), Inches(0.58),
             size=11, bold=True, color=ON_DARK, font="나눔고딕")

    # 오른쪽: Edge Loss + Freq Loss (11~12단계)
    add_rect(s, Inches(6.65), Inches(1.5), Inches(6.25), Inches(5.6), CARD)
    add_rect(s, Inches(6.65), Inches(1.5), Inches(6.25), Inches(0.1), TEAL)
    add_text(s, "11~12단계: 손실 함수 개선 (Edge + Frequency Loss)",
             Inches(6.85), Inches(1.72), Inches(5.85), Inches(0.42),
             size=13, bold=True, color=TEAL, font="나눔고딕")
    add_text(s,
             "목표: CNR SRAD(1.220) 초과\n"
             "현재 L1+SSIM은 경계 대비를 직접 최적화 못함\n\n"
             "Edge Loss (Sobel, 11단계):\n"
             "  정규화 6회 시도 끝에 SOBEL_MAX+clamp 확립\n"
             "  Loss += 0.5 * L1(clamp(Sobel/11.3), ...)\n\n"
             "Frequency Loss (FFT log1p, 12단계):\n"
             "  Loss += 0.1 * L1(log1p(|rfft2(pred)|), ...)\n"
             "  Edge + Freq 조합 시 오히려 CNR 하락",
             Inches(6.85), Inches(2.28), Inches(5.85), Inches(2.75),
             size=12, color=BODY, font="나눔고딕")
    loss_cols = ["방법", "PSNR", "SSIM", "CNR"]
    loss_rows = [
        ["NAFNet (8단계)", "28.11", "0.6743", "1.183"],
        ["+Edge L=0.5 (11단계)", "28.06", "0.6699", "1.186 (+)"],
        ["+Edge+Freq (12단계)", "28.05", "0.6683", "1.182 (-)"],
    ]
    buf2 = make_table_img(loss_cols, loss_rows, figsize=(6.2, 1.9),
                          highlight_rows={2}, fontsize=9, scale_y=2.2)
    s.shapes.add_picture(buf2, Inches(6.65), Inches(5.15), width=Inches(6.2))
    add_rect(s, Inches(6.65), Inches(6.38), Inches(6.25), Inches(0.72), DARK_ELV)
    add_text(s, "결론: Edge Loss로 CNR +0.003 확인. Freq 추가 시 gradient 충돌로 역효과. 손실 함수 개선 한계.",
             Inches(6.82), Inches(6.48), Inches(5.95), Inches(0.58),
             size=11, bold=True, color=ON_DARK, font="나눔고딕")


# ---------------------------------------------------------------------------
# 슬라이드 17 (NEW): 13단계 AROI N2N 사전학습
# ---------------------------------------------------------------------------

def slide_aroi_pretrain(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "13단계: AROI N2N 사전학습 → SBSDI D1 Fine-tuning")

    # 왼쪽: 방법론
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.6), CARD)
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(0.1), SUCCESS)
    add_text(s, "접근: 보유 데이터 최대 활용",
             Inches(0.6), Inches(1.72), Inches(5.5), Inches(0.42),
             size=14, bold=True, color=SUCCESS, font="나눔고딕")
    add_text(s,
             "AROI: 24명 × 128장 순차 B-scan = 3,072장\n"
             "인접 슬라이스 (i, i+1)는 조직 구조 유사 +\n"
             "스페클 패턴 독립 => N2N 조건 만족\n\n"
             "사전학습 (N2N):\n"
             "  Train: patient 1~20 (2,540 N2N 쌍)\n"
             "  Val: patient 21~24 (508 N2N 쌍)\n"
             "  Lazy loading + random crop (4,096/epoch)\n"
             "  NAFNet width=32 / L1 / lr=1e-3 / seed=42\n"
             "  Early stop epoch 93 | val loss 0.04305\n\n"
             "Fine-tuning (D1 k-fold):\n"
             "  사전학습 가중치 로드 후 lr=1e-4\n"
             "  Loss: L1+SSIM+EdgeLoss(λ=0.5)\n"
             "  모든 fold epoch 31~34에서 빠른 수렴",
             Inches(0.6), Inches(2.28), Inches(5.5), Inches(4.55),
             size=11.5, color=BODY, font="나눔고딕")

    # 오른쪽: D1 PSNR 모니터링 그래프 + 결과 표
    add_text(s, "사전학습 중 D1 PSNR 모니터링",
             Inches(6.65), Inches(1.55), Inches(6.2), Inches(0.42),
             size=13, bold=True, color=INK, font="나눔고딕")

    # D1 PSNR 모니터링 그래프 생성
    try:
        import pandas as pd
        df = pd.read_csv(str(PRETRAIN_LOG))
        d1 = df.dropna(subset=["d1_psnr"])
        fig, ax = plt.subplots(figsize=(6.0, 2.6))
        fig.patch.set_facecolor("#faf9f5")
        ax.plot(d1["epoch"], d1["d1_psnr"], color="#5db8a6", linewidth=2.0, marker="o", markersize=4)
        ax.axhline(27.50, color="#cc785c", linestyle="--", linewidth=1.5, label="SRAD (27.50)")
        ax.axhline(26.84, color="#e8a55a", linestyle=":", linewidth=1.5, label="N2N D2 (26.84)")
        ax.set_xlabel("Epoch", fontsize=9)
        ax.set_ylabel("D1 PSNR (dB)", fontsize=9)
        ax.set_title("AROI N2N Pre-training D1 PSNR", fontsize=10)
        ax.set_facecolor("#faf9f5")
        ax.grid(True, alpha=0.3, color="#e6dfd8")
        ax.legend(fontsize=8, framealpha=0)
        ax.tick_params(labelsize=8)
        for sp in ax.spines.values():
            sp.set_edgecolor("#e6dfd8")
        plt.tight_layout(pad=0.8)
        buf_graph = buf_to_slide_img(fig)
        s.shapes.add_picture(buf_graph, Inches(6.65), Inches(2.0), width=Inches(6.2))
    except Exception as e:
        print(f"  [WARN] D1 PSNR 그래프 생성 실패: {e}")

    res_cols = ["방법", "PSNR", "SSIM", "CNR"]
    res_rows = [
        ["SRAD",                "27.50", "0.652",  "1.220"],
        ["8단계 NAFNet",        "28.11", "0.6743", "1.183"],
        ["9단계 NAFNet+Aug",    "28.35", "0.6832", "1.168"],
        ["11단계 +Edge",        "28.06", "0.6699", "1.186"],
        ["13단계 AROI PT+FT",   "28.23", "0.6838*","1.177"],
    ]
    buf3 = make_table_img(res_cols, res_rows, figsize=(6.2, 2.4),
                          highlight_rows={5}, highlight_col=2, fontsize=9, scale_y=2.0)
    s.shapes.add_picture(buf3, Inches(6.65), Inches(4.72), width=Inches(6.2))
    add_text(s, "* SSIM 0.6838: 전체 방법 중 최고",
             Inches(6.65), Inches(6.88), Inches(6.2), Inches(0.38),
             size=11, bold=True, color=SUCCESS, font="나눔고딕")


# ---------------------------------------------------------------------------
# 슬라이드 19 (업데이트): 종합 성능 비교 차트 (16개 방법)
# ---------------------------------------------------------------------------

def slide_results_chart(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "종합 성능 비교 — 전 방법 PSNR / SSIM / CNR")

    labels = [
        "NLM", "BM3D", "SRAD",
        "N2N\n(3A)", "지도\n(3B)", "Pre+FT\n(5)",
        "ESR\nx2", "ESR\nx4",
        "U-Net\n(6)", "DnCNN\n(7)", "NAF\n(8)", "NAF\n+Aug",
        "lr1e-4\n(10)", "+Edge\n(11)", "+E+F\n(12)", "AROI\nPT(13)",
    ]
    psnr = [26.12, 27.00, 27.50, 26.84, 22.43, 27.46, 27.30, 27.57,
            28.19, 28.17, 28.11, 28.35, 27.95, 28.06, 28.05, 28.23]
    ssim = [0.492, 0.599, 0.652, 0.681, 0.333, 0.6795, 0.674, 0.673,
            0.6814, 0.6732, 0.6743, 0.6832, 0.6666, 0.6699, 0.6683, 0.6838]
    cnr  = [1.130, 1.134, 1.220, 1.148, 0.902, 1.169, 1.121, 1.166,
            1.169, 1.167, 1.183, 1.168, 1.191, 1.186, 1.182, 1.177]

    cat_colors = (
        ["#cc785c"] * 3 +
        ["#5db8a6", "#6c6a64", "#e8a55a"] +
        ["#5db8a6", "#5db872"] +
        ["#cc785c"] * 4 +
        ["#e8a55a"] * 4
    )

    fig, axes = plt.subplots(1, 3, figsize=(14.0, 4.0))
    fig.patch.set_facecolor("#faf9f5")
    data_sets = [
        (psnr, "PSNR (dB)", 27.50),
        (ssim, "SSIM",       0.652),
        (cnr,  "CNR",        1.220),
    ]
    for ax, (vals, label, ref) in zip(axes, data_sets):
        bars = ax.bar(labels, vals, color=cat_colors, width=0.62, zorder=3)
        ax.axhline(ref, color="#cc785c", linestyle="--", linewidth=1.5,
                   label=f"SRAD ({ref})", zorder=4)
        ax.set_facecolor("#faf9f5")
        ax.set_title(label, fontsize=11, color="#141413", fontweight="bold", pad=5)
        ax.tick_params(colors="#3d3d3a", labelsize=5.5)
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, rotation=40, ha="right", fontsize=5.5)
        for sp in ax.spines.values():
            sp.set_edgecolor("#e6dfd8")
        ax.yaxis.grid(True, color="#e6dfd8", zorder=0)
        ax.set_axisbelow(True)
        ymin = min(vals) * 0.94
        ymax = max(max(vals), ref) * 1.06
        ax.set_ylim(ymin, ymax)
        for bar, v in zip(bars, vals):
            ax.text(bar.get_x() + bar.get_width() / 2,
                    v + (ymax - ymin) * 0.008,
                    f"{v:.3f}", ha="center", va="bottom", fontsize=4.5, color="#141413")
        ax.legend(fontsize=7, framealpha=0)
    plt.tight_layout(pad=1.0)
    buf = buf_to_slide_img(fig)
    s.shapes.add_picture(buf, Inches(0.25), Inches(1.42), width=Inches(12.85))

    add_rect(s, Inches(0.25), Inches(6.62), Inches(12.85), Inches(0.72), DARK_ELV)
    add_text(s,
             "SSIM 최고: AROI pretrain(13단계) 0.6838  |  PSNR 최고: NAFNet+Aug(9단계) 28.35  |  CNR: AI 방법 전체 SRAD(1.220) 미달",
             Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.5),
             size=12, bold=True, color=ON_DARK, align=PP_ALIGN.CENTER, font="나눔고딕")


# ---------------------------------------------------------------------------
# 슬라이드 20 (업데이트): 종합 성능 비교 테이블 (16개 방법)
# ---------------------------------------------------------------------------

def slide_results_table(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "종합 성능 비교 테이블 (SBSDI D1, 18쌍 평균)")

    col_labels = ["방법", "PSNR (dB)", "SSIM", "CNR", "비고"]
    rows = [
        ["NLM",               "26.12+-2.01", "0.492+-0.050", "1.130+-0.119", "전통 방법"],
        ["BM3D",              "27.00+-2.51", "0.599+-0.066", "1.134+-0.119", "전통 방법"],
        ["SRAD",              "27.50+-1.98", "0.652+-0.023", "1.220+-0.121", "전통 기준값"],
        ["N2N (3A)",          "26.84+-2.18", "0.681+-0.033", "1.148+-0.133", "자가지도\nSSIM 초과"],
        ["지도학습 (3B)",     "22.43+-0.94", "0.333+-0.049", "0.902+-0.100", "도메인 갭\n크게 미달"],
        ["Pre-train+FT (5)",  "27.46+-2.55", "0.679+-0.037", "1.169+-0.120", "전이학습\nSRAD 근접"],
        ["ESRGAN x2 (4)",     "27.30+-2.35", "0.674+-0.034", "1.121+-0.116", "SR포함\nSSIM만 초과"],
        ["ESRGAN x4 (4)",     "27.57+-2.40", "0.673+-0.034", "1.166+-0.117", "SR포함\nPSNR+SSIM"],
        ["U-Net ES (6)",      "28.19+-2.56", "0.681+-0.031", "1.169+-0.121", "real GT\nSRAD 초과"],
        ["DnCNN (7)",         "28.17+-2.47", "0.673+-0.032", "1.167+-0.127", "from scratch\nU-Net 동등"],
        ["NAFNet-32 (8)",     "28.11+-2.26", "0.674+-0.029", "1.183+-0.120", "CNR AI 최고"],
        ["NAFNet+Aug (9)",    "28.35+-2.56", "0.683+-0.032", "1.168+-0.121", "PSNR 최고"],
        ["NAFNet lr1e-4 (10)","27.95+-2.11", "0.667+-0.026", "1.191+-0.116", "lr 낮추면 오히려\nPSNR 하락"],
        ["NAFNet+Edge (11)",  "28.06+-2.08", "0.670+-0.027", "1.186+-0.116", "CNR +0.003\nloss 튜닝 최선"],
        ["NAFNet+E+F (12)",   "28.05+-2.03", "0.668+-0.028", "1.182+-0.120", "Freq 추가\n역효과"],
        ["AROI PT+FT (13)",   "28.23+-2.49", "0.684+-0.029", "1.177+-0.121", "SSIM 전체\n최고"],
    ]
    buf = make_table_img(col_labels, rows, figsize=(13.5, 6.8),
                         highlight_rows={3, 5, 12, 16},
                         highlight_col=1,
                         fontsize=9, scale_y=2.0)
    s.shapes.add_picture(buf, Inches(0.3), Inches(1.42), width=Inches(12.8))

    add_text(s,
             "굵게: SRAD 기준(3행) / 도메인 갭(5행) / PSNR 최고(12행) / SSIM 최고(16행)  |  coral: PSNR 열 강조",
             Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.32),
             size=9, color=MUTED, font="나눔고딕")


# ---------------------------------------------------------------------------
# 슬라이드 21: 이미지 시각 비교 (v2와 동일)
# ---------------------------------------------------------------------------

def slide_image_comparison(prs):
    s = blank_slide(prs)
    set_bg(s, CANVAS)
    header_bar(s, "이미지 시각 비교 — SRAD vs DnCNN vs NAFNet+Aug (sample #01, fold 1)")
    labels = ["SRAD (기준값)", "DnCNN (7단계)", "NAFNet+Aug (9단계, 최고 PSNR)"]
    imgs   = [
        IMG_BASELINE / "sample_01.png",
        IMG_DNCNN    / "sample_01.png",
        IMG_AUG      / "sample_01.png",
    ]
    colors = [CORAL, AMBER, TEAL]
    for i, (lbl, img, col) in enumerate(zip(labels, imgs, colors)):
        cx = Inches(0.38 + i * 4.32)
        add_rect(s, cx, Inches(1.5), Inches(4.2), Inches(0.08), col)
        add_text(s, lbl, cx, Inches(1.62), Inches(4.2), Inches(0.42),
                 size=12, bold=True, color=INK, align=PP_ALIGN.CENTER, font="나눔고딕")
        add_img(s, img, cx, Inches(2.08), w=Inches(4.2))
    add_rect(s, Inches(0.38), Inches(6.78), Inches(12.94), Inches(0.62), DARK_ELV)
    add_text(s,
             "SRAD: noisy 단일 프레임 / clean 평균 / NLM / BM3D / SRAD  "
             "|  DnCNN & NAFNet+Aug: noisy / clean / 복원 (fold 1 평가 이미지)",
             Inches(0.55), Inches(6.88), Inches(12.6), Inches(0.45),
             size=10, color=ON_DARK_S, align=PP_ALIGN.CENTER, font="나눔고딕")


# ---------------------------------------------------------------------------
# 슬라이드 22 (업데이트): 핵심 발견 및 향후 계획
# ---------------------------------------------------------------------------

def slide_findings_future(prs):
    s = blank_slide(prs)
    set_bg(s, DARK)
    add_text(s, "핵심 발견 및 향후 계획", Inches(0.6), Inches(0.22),
             Inches(11), Inches(0.78), size=28, bold=False, color=ON_DARK, font="나눔고딕")
    add_divider(s, Inches(0.6), Inches(1.05), Inches(12.1), CORAL)

    add_rect(s, Inches(0.4), Inches(1.22), Inches(6.0), Inches(6.0), DARK_ELV)
    add_text(s, "핵심 발견 (13단계 기준)", Inches(0.6), Inches(1.42),
             Inches(5.6), Inches(0.42), size=14, bold=True, color=CORAL, font="나눔고딕")
    findings = [
        (CORAL,   "데이터가 주된 병목",
         "667K~17M 모델, lr 조정, 손실 함수 모두\nPSNR 28.0~28.4 dB 이내 — 아키텍처/HParam 차이 무의미"),
        (TEAL,    "clean GT가 결정적",
         "N2N(loss 0.184) => k-fold clean GT(loss 0.048)\n전환 시 PSNR 26.84 => 28.19 dB 도약"),
        (AMBER,   "사전학습 효과 확인",
         "AROI N2N 사전학습 => SSIM 0.6838 (전체 최고)\n실제 OCT 다양성으로 초기화가 SSIM 향상에 효과적"),
        (SUCCESS, "CNR 한계 지속",
         "어떤 AI 방법도 SRAD CNR(1.220) 초과 못함\n최고 1.191 — 경계 대비 보존이 남은 핵심 과제"),
    ]
    for i, (col, title, body) in enumerate(findings):
        add_rect(s, Inches(0.6), Inches(2.0 + i * 1.28), Inches(0.1), Inches(1.05), col)
        add_text(s, title, Inches(0.82), Inches(2.0 + i * 1.28),
                 Inches(5.4), Inches(0.42), size=12, bold=True, color=ON_DARK, font="나눔고딕")
        add_text(s, body, Inches(0.82), Inches(2.45 + i * 1.28),
                 Inches(5.4), Inches(0.72), size=10.5, color=ON_DARK_S, font="나눔고딕")

    add_rect(s, Inches(6.8), Inches(1.22), Inches(6.1), Inches(6.0), DARK_ELV)
    add_text(s, "향후 계획", Inches(7.0), Inches(1.42),
             Inches(5.7), Inches(0.42), size=14, bold=True, color=TEAL, font="나눔고딕")
    plans = [
        (CORAL,   "Duke AMD SD-OCT 활용",
         "38,400 B-scans 직접 다운로드 가능\n=> N2N 사전학습 데이터 다양성 대폭 확대"),
        (AMBER,   "외부 real OCT 데이터 확보",
         "PKU37 / Sub2Full vis-OCT 저자 컨택\nclean GT 확보로 데이터 병목 근본 해소"),
        (TEAL,    "치주 OCT 직접 수집",
         "동일 위치 ~40회 반복 스캔 => clean GT 생성\n교수님과 촬영 프로토콜 협의"),
        (SUCCESS, "스페클 제거 + SR 통합 파이프라인",
         "NAFNet 디노이징 후 Real-ESRGAN SR 연결\n치주 OCT 진단 품질 향상 최종 목표"),
    ]
    for i, (col, title, body) in enumerate(plans):
        add_rect(s, Inches(7.0), Inches(2.0 + i * 1.28), Inches(0.1), Inches(1.05), col)
        add_text(s, title, Inches(7.22), Inches(2.0 + i * 1.28),
                 Inches(5.5), Inches(0.42), size=12, bold=True, color=ON_DARK, font="나눔고딕")
        add_text(s, body, Inches(7.22), Inches(2.45 + i * 1.28),
                 Inches(5.5), Inches(0.72), size=10.5, color=ON_DARK_S, font="나눔고딕")

    add_divider(s, Inches(0.6), Inches(7.18), Inches(12.1), CORAL)
    add_text(s,
             "목표: 스페클 노이즈 제거 + SR 통합 AI 파이프라인 => 치주 OCT 임상 진단 품질 향상",
             Inches(0.6), Inches(7.25), Inches(12.1), Inches(0.22),
             size=12, bold=True, color=ON_DARK, align=PP_ALIGN.CENTER, font="나눔고딕")


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------

def main() -> None:
    prs = new_prs()
    print("슬라이드 생성 중...")

    slide_title(prs)
    print("  1/22 표지")

    slide_section(prs, "01", "OCT 기초 및 평가 지표",
                  ["OCT 원리: 근적외선 간섭계로 생체 조직 B-scan 재구성",
                   "스페클 노이즈: 레이저 가간섭성 => 무작위 후방 산란 패턴",
                   "노이즈 모델: I = S * Ns + Na  (Gamma 곱셈 + Gaussian 가산)",
                   "평가 지표: PSNR / SSIM (구조 유사도) / CNR (OCT 특화)"],
                  CORAL)
    print("  2/22 섹션1")

    slide_oct_basics(prs);    print("  3/22 OCT 원리")
    slide_metrics(prs);       print("  4/22 평가 지표")

    slide_section(prs, "02", "데이터",
                  ["SBSDI D1-D3: 인체/마우스 실제 OCT — clean GT는 D1(18쌍)만 존재",
                   "average.tif: 동일 위치 N~40회 반복 스캔 픽셀별 평균으로 생성",
                   "AROI 3,072장: 합성 1,136쌍 + 13단계 N2N 사전학습 3,048 인접쌍",
                   "Kermany 5,000쌍 합성 + 외부 탐색: 추가 GT 획득 불가"],
                  AMBER)
    print("  5/22 섹션2")

    slide_data_overview(prs); print("  6/22 데이터 현황")
    slide_data_survey(prs);   print("  7/22 외부 데이터 조사")

    slide_section(prs, "03", "방법론",
                  ["1~2단계: 전통 방법 기준값 + 물리 기반 합성 노이즈 파이프라인",
                   "3~5단계: N2N / 합성 지도학습(도메인 갭) / Pre-train+FT",
                   "6~9단계: 6-fold CV 아키텍처 비교 (U-Net/DnCNN/NAFNet) + 증강",
                   "10~12단계: lr 조정 + 손실 함수 개선 (Edge/Freq Loss)",
                   "13단계: AROI N2N 사전학습 => D1 fine-tuning (SSIM 최고)"],
                  TEAL)
    print("  8/22 섹션3")

    slide_pipeline_overview(prs); print("  9/22 파이프라인 개요")
    slide_traditional(prs);       print(" 10/22 1단계 전통 방법")
    slide_synthetic_noise(prs);   print(" 11/22 2단계 합성 노이즈")
    slide_n2n_supervised(prs);    print(" 12/22 3단계 N2N+지도학습")
    slide_sr_finetune(prs);       print(" 13/22 4~5단계 ESRGAN+FT")
    slide_kfold_archs(prs);       print(" 14/22 6~8단계 아키텍처 비교")
    slide_augmentation(prs);      print(" 15/22 9단계 증강")

    slide_lr_loss_exp(prs);       print(" 16/22 10~12단계 lr+손실함수")
    slide_aroi_pretrain(prs);     print(" 17/22 13단계 AROI 사전학습")

    slide_section(prs, "04", "결과 비교",
                  ["16개 방법 종합 PSNR / SSIM / CNR 비교 차트",
                   "전 방법 성능 비교 테이블 (SBSDI D1 18쌍 평균)",
                   "이미지 시각 비교: SRAD vs DnCNN vs NAFNet+Aug",
                   "핵심 발견 4가지 + 향후 계획"],
                  SUCCESS)
    print(" 18/22 섹션4")

    slide_results_chart(prs);      print(" 19/22 성능 비교 차트")
    slide_results_table(prs);      print(" 20/22 성능 비교 테이블")
    slide_image_comparison(prs);   print(" 21/22 이미지 시각 비교")
    slide_findings_future(prs);    print(" 22/22 핵심 발견 및 향후 계획")

    out = RESULTS_DIR / "presentation_v3.pptx"
    prs.save(str(out))
    print(f"\n저장 완료: {out}")


if __name__ == "__main__":
    main()
