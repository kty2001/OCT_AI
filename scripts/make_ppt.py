"""
OCT AI 프로젝트 발표 자료 자동 생성 스크립트.
결과: results/presentation.pptx

슬라이드 구성:
  1.  표지
  2.  연구 배경 및 목표
  3.  OCT와 스페클 노이즈
  4.  보유 데이터셋
  5.  평가 지표 (PSNR / SSIM / CNR)
  6.  방법론 전체 파이프라인
  7.  전통 방법 베이스라인 + 결과
  8.  합성 스페클 노이즈 생성 파이프라인 + 결과
  9.  3단계-A: 자가지도 학습 (N2N)
  10. 3단계-B: 지도학습 (합성 데이터 / 도메인 갭)
  11. 4단계: Real-ESRGAN SR 테스트
  12. 종합 성능 비교
  13. 분석 및 향후 계획
"""

import io
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm
import numpy as np
from PIL import Image

_korean_fonts = ["NanumGothic", "Malgun Gothic", "AppleGothic", "DejaVu Sans"]
for _fn in _korean_fonts:
    if any(_fn.lower() in f.name.lower() for f in fm.fontManager.ttflist):
        matplotlib.rcParams["font.family"] = _fn
        break

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ROOT        = Path(__file__).parent.parent
RESULTS_DIR = ROOT / "results"
IMG_BASELINE    = RESULTS_DIR / "01_baseline"   / "images"
IMG_SUB2FULL    = RESULTS_DIR / "03_sub2full"   / "images"
IMG_SUPERVISED  = RESULTS_DIR / "03_supervised" / "images"
IMG_SR_TEST     = RESULTS_DIR / "04_sr_test"
IMG_SYNTH       = RESULTS_DIR / "02_synthetic_noise"

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
CANVAS   = RGBColor(0xfa, 0xf9, 0xf5)
CORAL    = RGBColor(0xcc, 0x78, 0x5c)
CORAL_D  = RGBColor(0xa9, 0x58, 0x3e)
DARK     = RGBColor(0x18, 0x17, 0x15)
DARK_ELV = RGBColor(0x25, 0x23, 0x20)
CARD     = RGBColor(0xef, 0xe9, 0xde)
INK      = RGBColor(0x14, 0x14, 0x13)
BODY     = RGBColor(0x3d, 0x3d, 0x3a)
MUTED    = RGBColor(0x6c, 0x6a, 0x64)
ON_DARK  = RGBColor(0xfa, 0xf9, 0xf5)
ON_DARK_S= RGBColor(0xa0, 0x9d, 0x96)
AMBER    = RGBColor(0xe8, 0xa5, 0x5a)
TEAL     = RGBColor(0x5d, 0xb8, 0xa6)
SUCCESS  = RGBColor(0x5d, 0xb8, 0x72)
WHITE    = RGBColor(0xff, 0xff, 0xff)

SW = Inches(13.333)
SH = Inches(7.5)

SIZE_OFFSET = 3


# ---------------------------------------------------------------------------
# 유틸리티
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill


def add_text(slide, text: str, x, y, w, h,
             size: int = 18, bold: bool = False, color: RGBColor = INK,
             align=PP_ALIGN.LEFT, font: str = "나눔고딕") -> None:
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size + SIZE_OFFSET)
    run.font.bold  = bold
    run.font.color.rgb = color


def add_img(slide, path: Path, x, y, w=None, h=None) -> None:
    if not path.exists():
        print(f"  [WARN] 이미지 없음: {path}")
        return
    if w and h:
        slide.shapes.add_picture(str(path), x, y, w, h)
    elif w:
        slide.shapes.add_picture(str(path), x, y, width=w)
    elif h:
        slide.shapes.add_picture(str(path), x, y, height=h)
    else:
        slide.shapes.add_picture(str(path), x, y)


def add_divider(slide, x, y, w, color: RGBColor = CORAL) -> None:
    add_rect(slide, x, y, w, Pt(2), color)


def buf_to_slide_img(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=120)
    buf.seek(0)
    plt.close(fig)
    return buf


def make_table_img(col_labels, rows, figsize, highlight_row=None,
                   highlight_col=None, fontsize=10, scale_y=2.2):
    """matplotlib 테이블 이미지 생성."""
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor("#faf9f5")
    ax.axis("off")
    table = ax.table(cellText=rows, colLabels=col_labels,
                     cellLoc="center", loc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(fontsize)
    table.scale(1, scale_y)
    for (r, c), cell in table.get_celld().items():
        if r == 0:
            cell.set_facecolor("#181715")
            cell.get_text().set_color("#faf9f5")
            cell.get_text().set_fontweight("bold")
        elif highlight_row is not None and r == highlight_row:
            cell.set_facecolor("#e8ddd5")
            cell.get_text().set_fontweight("bold")
        elif r % 2 == 1:
            cell.set_facecolor("#efe9de")
        else:
            cell.set_facecolor("#faf9f5")
        if highlight_col is not None and c == highlight_col and r > 0:
            cell.get_text().set_color("#cc785c")
        cell.set_edgecolor("#e6dfd8")
    return buf_to_slide_img(fig)


# ---------------------------------------------------------------------------
# 슬라이드별 생성 함수
# ---------------------------------------------------------------------------

def slide_title(prs: Presentation) -> None:
    """슬라이드 1: 표지"""
    s = blank_slide(prs)
    set_bg(s, DARK)

    add_rect(s, Inches(0.6), Inches(1.5), Inches(0.06), Inches(4.0), CORAL)

    add_text(s, "AI 기반 OCT 스페클 노이즈 제거",
             Inches(0.85), Inches(1.6), Inches(9.5), Inches(1.6),
             size=40, bold=False, color=ON_DARK, font="나눔고딕")

    add_text(s, "망막 선행 연구를 기반으로 한 치주 도메인 전이 전략",
             Inches(0.85), Inches(3.3), Inches(9.0), Inches(0.8),
             size=22, bold=False, color=ON_DARK_S, font="나눔고딕")

    add_divider(s, Inches(0.85), Inches(4.3), Inches(4.0), CORAL)

    add_text(s, "OCT Speckle Denoising\nRetinal-to-Periodontal Transfer Strategy",
             Inches(0.85), Inches(4.5), Inches(8.0), Inches(0.8),
             size=14, color=ON_DARK_S, font="나눔고딕")

    add_rect(s, Inches(10.0), Inches(1.5), Inches(2.9), Inches(5.1), DARK_ELV)
    add_text(s, "진행 단계", Inches(10.2), Inches(1.7), Inches(2.5), Inches(0.4),
             size=11, color=CORAL, bold=True, font="나눔고딕")

    stages = [
        ("1단계  전통적 방법 베이스라인", SUCCESS),
        ("2단계  합성 노이즈 생성 파이프라인", SUCCESS),
        ("3단계-A  자가지도 학습 (N2N)", SUCCESS),
        ("3단계-B  지도학습 (합성 데이터)", SUCCESS),
        ("4단계  SR 사전 학습 모델 적용", CORAL),
    ]
    for i, (txt, col) in enumerate(stages):
        add_rect(s, Inches(10.2), Inches(2.1 + i * 0.72), Inches(0.1), Inches(0.4), col)
        add_text(s, txt, Inches(10.4), Inches(2.1 + i * 0.72), Inches(2.4), Inches(0.45),
                 size=11, color=ON_DARK if col != MUTED else ON_DARK_S, font="나눔고딕")


def slide_background(prs: Presentation) -> None:
    """슬라이드 2: 연구 배경 및 목표"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "연구 배경 및 목표", Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    cards = [
        ("문제 상황", CORAL,
         "OCT(광간섭단층촬영)는 비침습적으로 생체 조직 단면을\n"
         "마이크로미터 단위로 촬영하는 기술이다.\n"
         "레이저 가간섭성으로 스페클 노이즈가 필연적으로 발생하며\n"
         "이는 OCT 이미지 품질 저하와 진단 정확도 감소의 원인이다."),
        ("데이터 부족 문제", AMBER,
         "치주 OCT 공개 데이터는 거의 존재하지 않는다.\n"
         "clean Ground Truth 획득을 위한 다중 프레임 촬영은\n"
         "임상에서 촬영 시간이 크게 증가한다.\n"
         "대규모 안과 망막 OCT 데이터 활용이 필수적이다."),
        ("연구 목표", TEAL,
         "1. 전통적 방법(NLM, BM3D, SRAD)으로 기준값 확보\n"
         "2. 물리 기반 노이즈 모델로 합성 학습 데이터 생성\n"
         "3. clean GT 없이 자가지도/지도 학습 디노이징 구현\n"
         "4. 사전 학습 SR 모델로 노이즈 제거 + 해상도 향상"),
    ]
    for i, (title, col, body) in enumerate(cards):
        cx = Inches(0.4 + i * 4.28)
        add_rect(s, cx, Inches(1.6), Inches(4.1), Inches(5.5), CARD)
        add_rect(s, cx, Inches(1.6), Inches(4.1), Inches(0.08), col)
        add_text(s, title, cx + Inches(0.2), Inches(1.78), Inches(3.7), Inches(0.5),
                 size=15, bold=True, color=INK, font="나눔고딕")
        add_text(s, body, cx + Inches(0.2), Inches(2.4), Inches(3.7), Inches(4.4),
                 size=12, color=BODY, font="나눔고딕")


def slide_oct_noise(prs: Presentation) -> None:
    """슬라이드 3: OCT와 스페클 노이즈"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "OCT와 스페클 노이즈", Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    add_rect(s, Inches(0.4), Inches(1.6), Inches(5.5), Inches(5.5), CARD)
    add_text(s, "OCT 원리", Inches(0.6), Inches(1.8), Inches(5.0), Inches(0.5),
             size=16, bold=True, color=CORAL, font="나눔고딕")
    oct_text = (
        "근적외선 마이켈슨 간섭계로 생체 조직의 깊이 방향\n"
        "단면(B-scan)을 재구성한다.\n\n"
        "A-scan: 깊이 방향 1D 신호\n"
        "B-scan: A-scan 연속 배열로 구성된 2D 단면\n\n"
        "스페클 노이즈 원인:\n"
        "  레이저 가간섭성 -> 후방 산란 위상 간섭\n"
        "  촬영마다 무작위로 다른 스페클 패턴 발생\n\n"
        "노이즈 모델 (곱셈성 + 가산성):\n"
        "  I = S * Ns + Na\n"
        "  Ns ~ Gamma(L, 1/L)  [곱셈성 스페클]\n"
        "  Na ~ N(0, sigma^2)  [가산성 가우시안]\n\n"
        "SBSDI D1 캘리브레이션:\n"
        "  L = 5.27,  sigma = 0.010"
    )
    add_text(s, oct_text, Inches(0.6), Inches(2.4), Inches(5.1), Inches(4.5),
             size=12, color=BODY, font="나눔고딕")

    add_text(s, "실제 OCT B-scan 예시 (SBSDI D1)", Inches(6.3), Inches(1.65), Inches(6.6), Inches(0.4),
             size=13, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_BASELINE / "sample_01.png", Inches(6.1), Inches(2.1), w=Inches(6.9))
    add_text(s, "왼쪽: 단일 프레임 noisy B-scan  /  오른쪽: 다중 프레임 평균 clean",
             Inches(6.1), Inches(6.75), Inches(6.9), Inches(0.4),
             size=10, color=MUTED, font="나눔고딕")


def slide_dataset(prs: Presentation) -> None:
    """슬라이드 4: 보유 데이터셋"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "보유 데이터셋", Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    col_labels = ["데이터셋", "이미지 수", "태스크", "noisy-clean", "활용"]
    rows = [
        ["SBSDI D1\n(Fang 2013)", "18쌍 + 39세트", "디노이징", "있음 (18쌍)", "기준값 / N2N 학습"],
        ["AROI\n(Heidelberg)", "1,136장 주석", "레이어 세그", "없음 -> 합성", "합성 학습 데이터"],
        ["Kermany OCT2017\n(CC BY 4.0)", "84,484장", "4클래스 분류", "없음 -> 합성", "합성 학습 데이터"],
    ]
    buf = make_table_img(col_labels, rows, figsize=(9.5, 3.2), fontsize=10, scale_y=2.3)
    s.shapes.add_picture(buf, Inches(0.4), Inches(1.5), width=Inches(7.8))

    add_text(s, "AROI 합성 노이즈 샘플", Inches(8.6), Inches(1.55), Inches(4.3), Inches(0.4),
             size=13, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_SYNTH / "AROI" / "samples" / "sample_00000.png",
            Inches(8.5), Inches(2.0), w=Inches(4.5))
    add_text(s, "위: clean  /  아래: 합성 스페클 노이즈 적용",
             Inches(8.5), Inches(6.6), Inches(4.5), Inches(0.4),
             size=10, color=MUTED, font="나눔고딕")

    add_rect(s, Inches(0.4), Inches(6.05), Inches(7.8), Inches(1.05), DARK_ELV)
    add_text(s,
             "합성 노이즈 생성 결과: AROI 1,136쌍 + Kermany 5,000쌍 = 총 6,136 noisy-clean 쌍",
             Inches(0.6), Inches(6.2), Inches(7.4), Inches(0.7),
             size=12, bold=True, color=ON_DARK, font="나눔고딕")


def slide_metrics(prs: Presentation) -> None:
    """슬라이드 5: 평가 지표"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "평가 지표", Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    metrics = [
        ("PSNR", "Peak Signal-to-Noise Ratio\n최대 신호 대 잡음비", CORAL,
         "PSNR = 10 * log10( MAX^2 / MSE )",
         "MSE: 복원/clean 픽셀별 제곱 오차 평균\n"
         "MAX: 픽셀 최댓값 (0-1 기준 1.0)\n\n"
         "오차가 작을수록 PSNR 높음 => 클수록 좋음\n"
         "단위: dB, 실용 범위 20-40 dB\n\n"
         "한계: 픽셀 수치 오차만 반영"),
        ("SSIM", "Structural Similarity Index\n구조적 유사도", AMBER,
         "SSIM = f(밝기, 대비, 구조 공분산)",
         "밝기(mu) + 대비(sigma^2) + 구조(sigma_xy)\n"
         "동시 비교 => 인간 시각에 가까운 평가\n\n"
         "범위: 0-1, 높을수록 좋음\n"
         "PSNR보다 임상 판독 품질에 가까운 지표\n\n"
         "한계: 지역 윈도(11x11) 기반"),
        ("CNR", "Contrast-to-Noise Ratio\n대비 대 잡음비", TEAL,
         "CNR = |mu_signal - mu_bg| / sigma_bg",
         "OCT 조직 레이어(신호)와 배경(잡음) 간\n"
         "대비를 잡음 수준으로 나눈 값\n\n"
         "높을수록 레이어 경계 선명, 구조 식별 용이\n\n"
         "임상적 진단 품질에 직접 대응하는\n"
         "OCT 특화 지표"),
    ]

    for i, (abbr, full, col, formula, desc) in enumerate(metrics):
        cx = Inches(0.4 + i * 4.28)
        add_rect(s, cx, Inches(1.55), Inches(4.1), Inches(5.55), CARD)
        add_rect(s, cx, Inches(1.55), Inches(4.1), Inches(0.1), col)
        add_text(s, abbr, cx + Inches(0.2), Inches(1.75), Inches(3.7), Inches(0.6),
                 size=22, bold=True, color=INK, font="나눔고딕")
        add_text(s, full, cx + Inches(0.2), Inches(2.4), Inches(3.7), Inches(0.55),
                 size=11, color=MUTED, font="나눔고딕")
        add_divider(s, cx + Inches(0.2), Inches(3.0), Inches(3.5), col)
        add_rect(s, cx + Inches(0.2), Inches(3.1), Inches(3.7), Inches(0.6), DARK_ELV)
        add_text(s, formula, cx + Inches(0.3), Inches(3.17), Inches(3.5), Inches(0.5),
                 size=10, color=ON_DARK, font="Courier New")
        add_text(s, desc, cx + Inches(0.2), Inches(3.8), Inches(3.7), Inches(3.15),
                 size=11, color=BODY, font="나눔고딕")


def slide_pipeline(prs: Presentation) -> None:
    """슬라이드 6: 방법론 전체 파이프라인"""
    s = blank_slide(prs)
    set_bg(s, DARK)

    add_text(s, "연구 방법론 전체 파이프라인", Inches(0.6), Inches(0.25), Inches(11), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")
    add_divider(s, Inches(0.6), Inches(1.1), Inches(12.0), CORAL)

    stages = [
        ("1단계", "전통적 방법\n베이스라인", CORAL,
         "NLM / BM3D / SRAD\nSBSDI D1 18쌍 평가",
         "SRAD: PSNR 27.50\nSSIM 0.652"),
        ("2단계", "합성 노이즈\n생성 파이프라인", AMBER,
         "I = S*Ns + Na\nGamma 캘리브레이션\n6,136쌍 생성",
         "L=5.27, KS=0.119"),
        ("3단계", "자가지도 +\n지도학습", TEAL,
         "N2N 프레임쌍 (3-A)\n합성 지도학습 (3-B)\nU-Net ~2M params",
         "N2N: SSIM 0.681\n지도: 도메인 갭 확인"),
        ("4단계", "사전 학습\nSR 모델 적용", CORAL,
         "Real-ESRGAN\n(blind SR + 노이즈 제거)\nx2 / x4 업스케일",
         "PSNR 27.92 > SRAD\nSSIM 0.675 > SRAD"),
    ]

    for i, (num, title, col, body, result) in enumerate(stages):
        cx = Inches(0.35 + i * 3.2)
        bg_col = DARK_ELV
        add_rect(s, cx, Inches(1.3), Inches(3.0), Inches(5.8), bg_col)
        add_rect(s, cx, Inches(1.3), Inches(3.0), Inches(0.12), col)

        add_text(s, num, cx + Inches(0.15), Inches(1.5), Inches(2.7), Inches(0.38),
                 size=11, bold=True, color=col, font="나눔고딕")
        add_text(s, title, cx + Inches(0.15), Inches(1.88), Inches(2.7), Inches(0.75),
                 size=15, bold=True, color=ON_DARK, font="나눔고딕")
        add_divider(s, cx + Inches(0.15), Inches(2.65), Inches(2.65), col)
        add_text(s, body, cx + Inches(0.15), Inches(2.8), Inches(2.7), Inches(2.4),
                 size=12, color=ON_DARK, font="나눔고딕")

        add_rect(s, cx + Inches(0.15), Inches(5.45), Inches(2.65), Inches(0.07), col)
        add_text(s, "결과", cx + Inches(0.15), Inches(5.58), Inches(0.7), Inches(0.3),
                 size=10, bold=True, color=col, font="나눔고딕")
        add_text(s, result, cx + Inches(0.15), Inches(5.88), Inches(2.65), Inches(0.9),
                 size=11, color=ON_DARK, font="나눔고딕")

        if i < 3:
            add_text(s, ">", cx + Inches(3.0), Inches(3.8), Inches(0.2), Inches(0.5),
                     size=20, bold=True, color=MUTED, font="나눔고딕", align=PP_ALIGN.CENTER)


def slide_baseline(prs: Presentation) -> None:
    """슬라이드 7: 전통 방법 베이스라인"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "1단계: 전통적 방법 베이스라인", Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    methods = [
        ("NLM", "Non-Local Means", CORAL,
         "이미지 전체에서 유사 패치를 찾아\n가중 평균으로 노이즈 제거\n속도: ~0.5초/장"),
        ("BM3D", "Block-Matching 3D", AMBER,
         "유사 블록을 3D 변환 도메인으로\n묶어 임계화 후 재구성\n속도: ~3.2초/장"),
        ("SRAD", "Speckle Reducing\nAnisotropic Diffusion", TEAL,
         "곱셈성 스페클 모델 반영 PDE 기반\n이방성 확산, 경계 보존 최우수\n속도: ~5.7초/장"),
    ]
    for i, (abbr, full, col, desc) in enumerate(methods):
        cx = Inches(0.4 + i * 3.1)
        add_rect(s, cx, Inches(1.5), Inches(2.9), Inches(2.4), CARD)
        add_rect(s, cx, Inches(1.5), Inches(2.9), Inches(0.08), col)
        add_text(s, abbr, cx + Inches(0.15), Inches(1.65), Inches(2.6), Inches(0.48),
                 size=16, bold=True, color=INK, font="나눔고딕")
        add_text(s, full, cx + Inches(0.15), Inches(2.15), Inches(2.6), Inches(0.45),
                 size=10, color=MUTED, font="나눔고딕")
        add_text(s, desc, cx + Inches(0.15), Inches(2.65), Inches(2.6), Inches(1.15),
                 size=11, color=BODY, font="나눔고딕")

    col_labels = ["방법", "PSNR (dB)", "SSIM", "CNR"]
    rows = [
        ["NLM",  "26.12 +- 2.01", "0.492 +- 0.050", "1.130"],
        ["BM3D", "27.00 +- 2.51", "0.599 +- 0.066", "1.134"],
        ["SRAD", "27.50 +- 1.98", "0.652 +- 0.023", "1.220"],
    ]
    buf = make_table_img(col_labels, rows, figsize=(5.0, 2.0),
                         highlight_row=3, fontsize=10, scale_y=2.2)
    s.shapes.add_picture(buf, Inches(0.4), Inches(4.05), width=Inches(5.0))

    add_text(s, "SRAD가 PSNR / SSIM / CNR 전 지표 1위\n딥러닝 목표: PSNR > 27.50, SSIM > 0.652",
             Inches(0.4), Inches(6.45), Inches(5.0), Inches(0.8),
             size=11, bold=True, color=CORAL_D, font="나눔고딕")

    add_text(s, "시각적 비교 (sample #01: Noisy / Clean / NLM / BM3D / SRAD)",
             Inches(5.7), Inches(3.85), Inches(7.3), Inches(0.45),
             size=11, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_BASELINE / "sample_01.png", Inches(5.7), Inches(4.35), w=Inches(7.3))


def slide_synthetic(prs: Presentation) -> None:
    """슬라이드 8: 합성 노이즈 생성 파이프라인"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "2단계: 합성 스페클 노이즈 생성 파이프라인", Inches(0.6), Inches(0.3), Inches(11), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    add_rect(s, Inches(0.4), Inches(1.55), Inches(5.3), Inches(5.55), CARD)
    add_text(s, "물리 기반 노이즈 모델", Inches(0.6), Inches(1.75), Inches(5.0), Inches(0.42),
             size=14, bold=True, color=CORAL, font="나눔고딕")

    add_rect(s, Inches(0.6), Inches(2.22), Inches(4.9), Inches(0.75), DARK_ELV)
    add_text(s, "I = S * Ns + Na\nNs ~ Gamma(L, 1/L)    Na ~ N(0, sigma^2)",
             Inches(0.7), Inches(2.28), Inches(4.7), Inches(0.65),
             size=11, color=ON_DARK, font="Courier New")

    model_desc = (
        "I : 합성 noisy 이미지\n"
        "S : 원본 clean 이미지\n"
        "Ns: 곱셈성 스페클  (mean=1, var=1/L)\n"
        "Na: 가산성 가우시안 노이즈\n"
        "L : looks 수 — 높을수록 노이즈 약함"
    )
    add_text(s, model_desc, Inches(0.6), Inches(3.1), Inches(4.9), Inches(1.45),
             size=12, color=BODY, font="나눔고딕")

    add_divider(s, Inches(0.6), Inches(4.65), Inches(4.9), AMBER)
    add_text(s, "파라미터 캘리브레이션 (SBSDI D1 18쌍)",
             Inches(0.6), Inches(4.75), Inches(4.9), Inches(0.42),
             size=12, bold=True, color=INK, font="나눔고딕")
    calib_text = (
        "Ns = I / S 의 Gamma 모멘트 매칭\n\n"
        "  L     = 5.266 +- 0.743  [3.015, 6.228]\n"
        "  sigma = 0.010  (모든 쌍 일정)\n"
        "  KS 통계량 = 0.119"
    )
    add_text(s, calib_text, Inches(0.6), Inches(5.22), Inches(4.9), Inches(1.65),
             size=12, color=BODY, font="나눔고딕")

    add_text(s, "합성 노이즈 생성 샘플 (AROI)", Inches(6.0), Inches(1.6), Inches(7.0), Inches(0.42),
             size=13, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_SYNTH / "AROI" / "samples" / "sample_00000.png",
            Inches(6.0), Inches(2.08), w=Inches(7.0))

    add_rect(s, Inches(6.0), Inches(5.65), Inches(7.0), Inches(1.45), DARK_ELV)
    add_text(s, "생성된 학습 데이터", Inches(6.2), Inches(5.78), Inches(6.6), Inches(0.38),
             size=11, bold=True, color=CORAL, font="나눔고딕")
    summary = (
        "AROI         1,136쌍   512x1024px   90도 회전 보정\n"
        "Kermany  5,000쌍   496x512px    랜덤 샘플링 (seed=42)\n"
        "합계          6,136쌍   재현 가능 파이프라인"
    )
    add_text(s, summary, Inches(6.2), Inches(6.2), Inches(6.6), Inches(0.85),
             size=11, color=ON_DARK, font="나눔고딕")


def slide_n2n(prs: Presentation) -> None:
    """슬라이드 9: 3단계-A 자가지도 학습 (N2N)"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "3단계-A: 자가지도 학습 (Noise2Noise)", Inches(0.6), Inches(0.3), Inches(11), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    add_rect(s, Inches(0.4), Inches(1.55), Inches(5.8), Inches(5.55), CARD)
    add_text(s, "핵심 아이디어", Inches(0.6), Inches(1.75), Inches(5.4), Inches(0.42),
             size=14, bold=True, color=TEAL, font="나눔고딕")
    idea = (
        "clean Ground Truth 없이 OCT 이미지만으로 학습.\n\n"
        "같은 위치를 반복 스캔한 프레임(1~4.tif)은\n"
        "조직 구조는 동일하고 스페클은 독립적이다.\n\n"
        "(frame_i, frame_j) 쌍 -> noisy-noisy 학습 쌍\n"
        "N2N 이론: 두 독립 노이즈 이미지로 학습 시\n"
        "최소 분산 추정 = clean 이미지에 수렴"
    )
    add_text(s, idea, Inches(0.6), Inches(2.28), Inches(5.4), Inches(2.25),
             size=12, color=BODY, font="나눔고딕")

    add_divider(s, Inches(0.6), Inches(4.6), Inches(5.4), TEAL)
    add_text(s, "학습 설정", Inches(0.6), Inches(4.75), Inches(5.4), Inches(0.42),
             size=13, bold=True, color=INK, font="나눔고딕")
    config = (
        "학습 데이터: 39세트 x 12 순서쌍 = 468쌍\n"
        "아키텍처: 경량 U-Net  (~1.95M params)\n"
        "손실 함수: L1 Loss  /  Adam lr=1e-4 (CosineAnnealing)\n"
        "에포크: 500  /  배치: 4  /  패치: 128x128\n"
        "학습 시간: ~221분 (RTX A4000)"
    )
    add_text(s, config, Inches(0.6), Inches(5.22), Inches(5.4), Inches(1.65),
             size=12, color=BODY, font="나눔고딕")

    add_text(s, "복원 결과 비교 (sample #01)",
             Inches(6.5), Inches(1.6), Inches(6.5), Inches(0.42),
             size=13, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_SUB2FULL / "sample_01.png", Inches(6.5), Inches(2.08), w=Inches(6.5))

    add_rect(s, Inches(6.5), Inches(5.65), Inches(6.5), Inches(1.45), DARK_ELV)
    add_text(s, "N2N 결과 (SBSDI D1 18쌍 평균)",
             Inches(6.7), Inches(5.78), Inches(6.1), Inches(0.38),
             size=11, bold=True, color=TEAL, font="나눔고딕")
    perf = (
        "PSNR  26.84 dB +- 2.18   (SRAD: 27.50)  미달\n"
        "SSIM   0.681 +- 0.033    (SRAD:  0.652)  [초과]\n"
        "CNR    1.148 +- 0.133    (SRAD:  1.220)  미달"
    )
    add_text(s, perf, Inches(6.7), Inches(6.22), Inches(6.1), Inches(0.82),
             size=11, color=ON_DARK, font="Courier New")


def slide_supervised(prs: Presentation) -> None:
    """슬라이드 10: 3단계-B 지도학습 (합성 데이터 / 도메인 갭)"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "3단계-B: 지도학습 — 합성 데이터 (도메인 갭 확인)",
             Inches(0.6), Inches(0.3), Inches(11.5), Inches(0.8),
             size=28, bold=False, color=ON_DARK, font="나눔고딕")

    # 왼쪽: 학습 설정 + 결과
    add_rect(s, Inches(0.4), Inches(1.55), Inches(5.8), Inches(5.55), CARD)
    add_text(s, "학습 설정", Inches(0.6), Inches(1.72), Inches(5.4), Inches(0.42),
             size=14, bold=True, color=CORAL, font="나눔고딕")
    config = (
        "학습 데이터: 6,136쌍 (AROI 1,136 + Kermany 5,000)\n"
        "아키텍처: U-Net (~1.95M params, N2N 모델 공유)\n"
        "손실 함수: L1 Loss  /  Adam lr=2e-4 (CosineAnnealing)\n"
        "에포크: 100  /  배치: 8  /  패치: 128x128\n"
        "학습 시간: ~305분 (RTX A4000)"
    )
    add_text(s, config, Inches(0.6), Inches(2.22), Inches(5.4), Inches(1.65),
             size=12, color=BODY, font="나눔고딕")

    add_divider(s, Inches(0.6), Inches(3.95), Inches(5.4), CORAL)
    add_text(s, "Loss 수렴 추이", Inches(0.6), Inches(4.08), Inches(5.4), Inches(0.4),
             size=12, bold=True, color=INK, font="나눔고딕")
    loss_text = (
        "Epoch   1:  0.044891\n"
        "Epoch  10:  0.023456   <- epoch 10 이후 수렴 포화\n"
        "Epoch  50:  0.021973\n"
        "Epoch 100:  0.021655   best: 0.021621"
    )
    add_text(s, loss_text, Inches(0.6), Inches(4.55), Inches(5.4), Inches(1.25),
             size=11, color=BODY, font="Courier New")

    add_divider(s, Inches(0.6), Inches(5.88), Inches(5.4), CORAL)
    add_text(s, "결과 (SBSDI D1 18쌍 평균)", Inches(0.6), Inches(6.0), Inches(5.4), Inches(0.4),
             size=12, bold=True, color=INK, font="나눔고딕")
    result_text = (
        "PSNR  22.43 dB  (SRAD: 27.50)  미달\n"
        "SSIM   0.333    (SRAD:  0.652)  미달\n"
        "CNR    0.902    (SRAD:  1.220)  미달"
    )
    add_text(s, result_text, Inches(0.6), Inches(6.48), Inches(5.4), Inches(0.85),
             size=11, color=CORAL_D, font="Courier New")

    # 오른쪽: 도메인 갭 분석 + 샘플 이미지
    add_text(s, "도메인 갭 (Domain Gap)", Inches(6.5), Inches(1.6), Inches(6.5), Inches(0.42),
             size=14, bold=True, color=AMBER, font="나눔고딕")

    add_rect(s, Inches(6.5), Inches(2.1), Inches(6.5), Inches(2.1), DARK_ELV)
    gap_col_labels = ["구분", "학습 데이터", "평가 데이터"]
    gap_rows = [
        ["노이즈 유형", "합성 Gamma+Gaussian", "실제 OCT 스페클"],
        ["이미지 소스", "AROI / Kermany", "SBSDI D1"],
        ["Clean 기준", "수학적 모델 생성", "다중 프레임 평균"],
    ]
    buf = make_table_img(gap_col_labels, gap_rows,
                         figsize=(6.8, 2.0), fontsize=10, scale_y=2.0)
    s.shapes.add_picture(buf, Inches(6.5), Inches(2.1), width=Inches(6.5))

    add_text(s, "비교: N2N (실제 데이터 학습) vs 지도학습 (합성 데이터)",
             Inches(6.5), Inches(4.38), Inches(6.5), Inches(0.45),
             size=12, bold=True, color=INK, font="나눔고딕")

    comp_col_labels = ["방법", "PSNR", "SSIM", "CNR"]
    comp_rows = [
        ["SRAD (기준)", "27.50", "0.652", "1.220"],
        ["N2N 프레임쌍", "26.84", "0.681 *", "1.148"],
        ["지도학습 (합성)", "22.43", "0.333", "0.902"],
    ]
    buf2 = make_table_img(comp_col_labels, comp_rows,
                          figsize=(6.8, 2.2), fontsize=10, scale_y=2.0)
    s.shapes.add_picture(buf2, Inches(6.5), Inches(4.9), width=Inches(6.5))

    add_text(s, "* SRAD 초과  |  합성 데이터 노이즈 모델 정확도가 지도학습 성패를 좌우함",
             Inches(6.5), Inches(6.82), Inches(6.5), Inches(0.45),
             size=10, color=MUTED, font="나눔고딕")


def slide_sr_test(prs: Presentation) -> None:
    """슬라이드 11: 4단계 Real-ESRGAN SR 테스트"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "4단계: 사전 학습 SR 모델 — Real-ESRGAN",
             Inches(0.6), Inches(0.3), Inches(11.5), Inches(0.8),
             size=28, bold=False, color=ON_DARK, font="나눔고딕")

    # 왼쪽: 모델 설명 + 결과 표
    add_rect(s, Inches(0.4), Inches(1.55), Inches(5.8), Inches(5.55), CARD)
    add_text(s, "Real-ESRGAN 모델", Inches(0.6), Inches(1.72), Inches(5.4), Inches(0.42),
             size=14, bold=True, color=TEAL, font="나눔고딕")
    model_desc = (
        "Blind SR: 노이즈 + 블러 + 압축 아티팩트 등\n"
        "실제 복합 열화를 동시에 복원하는 사전 학습 모델\n\n"
        "아키텍처: RRDBNet (잔차 밀집 블록)\n"
        "가중치: RealESRGAN_x4plus.pth / x2plus.pth (64MB)\n"
        "추론: CUDA half precision, tile=400"
    )
    add_text(s, model_desc, Inches(0.6), Inches(2.22), Inches(5.4), Inches(1.75),
             size=12, color=BODY, font="나눔고딕")

    add_divider(s, Inches(0.6), Inches(4.05), Inches(5.4), TEAL)
    add_text(s, "테스트 결과 (SBSDI D1 세트 1, 원본 스케일 비교)",
             Inches(0.6), Inches(4.18), Inches(5.4), Inches(0.42),
             size=12, bold=True, color=INK, font="나눔고딕")

    col_labels = ["방법", "PSNR (dB)", "SSIM", "출력 크기", "시간"]
    rows = [
        ["SRAD (기준)",      "27.50", "0.652", "450x900",   "5.7s"],
        ["Real-ESRGAN x2",  "27.69", "0.673", "900x1800",  "0.93s"],
        ["Real-ESRGAN x4",  "27.92", "0.675", "1800x3600", "1.47s"],
    ]
    buf = make_table_img(col_labels, rows, figsize=(5.6, 2.0),
                         highlight_row=None, fontsize=10, scale_y=2.2)
    s.shapes.add_picture(buf, Inches(0.4), Inches(4.65), width=Inches(5.8))

    add_text(s, "* 지표는 SR 출력을 원본 해상도로 다운스케일 후 clean과 비교",
             Inches(0.4), Inches(6.68), Inches(5.8), Inches(0.45),
             size=10, color=MUTED, font="나눔고딕")

    # 오른쪽: 비교 이미지
    add_text(s, "비교 이미지 (노이즈 입력 / 클린 레퍼런스 / x2 SR / x4 SR)",
             Inches(6.5), Inches(1.6), Inches(6.5), Inches(0.45),
             size=12, bold=True, color=INK, font="나눔고딕")
    add_img(s, IMG_SR_TEST / "all_comparison.png", Inches(6.5), Inches(2.1), w=Inches(6.5))

    add_rect(s, Inches(6.5), Inches(6.55), Inches(6.5), Inches(0.75), DARK_ELV)
    add_text(s, "PSNR / SSIM 모두 SRAD(전통 방법 1위) 초과  |  처리 속도 4~6배 향상",
             Inches(6.7), Inches(6.65), Inches(6.1), Inches(0.55),
             size=12, bold=True, color=ON_DARK, font="나눔고딕")


def slide_comparison(prs: Presentation) -> None:
    """슬라이드 12: 종합 성능 비교"""
    s = blank_slide(prs)
    set_bg(s, CANVAS)

    add_rect(s, Inches(0), Inches(0), SW, Inches(1.4), DARK)
    add_text(s, "종합 성능 비교", Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")

    methods   = ["NLM", "BM3D", "SRAD", "N2N\n(3-A)", "Real-ESRGAN\nx4 (4단계)"]
    psnr_vals = [26.12, 27.00, 27.50, 26.84, 27.92]
    ssim_vals = [0.492, 0.599, 0.652, 0.681, 0.675]
    cnr_vals  = [1.130, 1.134, 1.220, 1.148, None]

    colors_bar = ["#a09d96", "#a09d96", "#cc785c", "#5db8a6", "#5db872"]

    fig, axes = plt.subplots(1, 3, figsize=(10.5, 3.6))
    fig.patch.set_facecolor("#faf9f5")

    for ax, vals_raw, label, ref_val in zip(
        axes,
        [psnr_vals, ssim_vals, cnr_vals],
        ["PSNR (dB)", "SSIM", "CNR"],
        [27.50,       0.652,  1.220],
    ):
        vals_plot = [v if v is not None else 0 for v in vals_raw]
        colors_use = [c if (vals_raw[i] is not None) else "#3d3d3a"
                      for i, c in enumerate(colors_bar)]
        bars = ax.bar(methods, vals_plot, color=colors_use, width=0.55, zorder=3)
        ax.axhline(ref_val, color="#cc785c", linestyle="--", linewidth=1.5,
                   label=f"SRAD ({ref_val})", zorder=4)
        ax.set_facecolor("#faf9f5")
        ax.set_title(label, fontsize=11, color="#141413", fontweight="bold", pad=5)
        ax.tick_params(colors="#3d3d3a", labelsize=8)
        for spine in ax.spines.values():
            spine.set_edgecolor("#e6dfd8")
        ax.yaxis.grid(True, color="#e6dfd8", zorder=0)
        ax.set_axisbelow(True)
        valid = [v for v in vals_raw if v is not None]
        ymin = min(valid) * 0.95
        ymax = max(max(valid), ref_val) * 1.06
        ax.set_ylim(ymin, ymax)
        for bar, v in zip(bars, vals_raw):
            if v is not None:
                ax.text(bar.get_x() + bar.get_width() / 2, v + (ymax - ymin) * 0.01,
                        f"{v:.3f}", ha="center", va="bottom", fontsize=7.5, color="#141413")
            else:
                ax.text(bar.get_x() + bar.get_width() / 2, ymin + (ymax - ymin) * 0.05,
                        "N/A", ha="center", va="bottom", fontsize=7.5, color="#6c6a64")
        ax.legend(fontsize=8, framealpha=0)

    plt.tight_layout(pad=1.5)
    buf = buf_to_slide_img(fig)
    s.shapes.add_picture(buf, Inches(0.3), Inches(1.5), width=Inches(9.2))

    add_rect(s, Inches(9.6), Inches(1.5), Inches(3.4), Inches(5.6), DARK_ELV)
    add_text(s, "핵심 발견", Inches(9.8), Inches(1.65), Inches(3.0), Inches(0.42),
             size=13, bold=True, color=CORAL, font="나눔고딕")

    findings = [
        (SUCCESS, "Real-ESRGAN x4 최고 PSNR",
         "27.92 dB > SRAD 27.50\nSR 사전 학습으로 노이즈+SR 동시 해결"),
        (TEAL,    "N2N 최고 SSIM",
         "0.681 > SRAD 0.652\nclean GT 없이 구조 보존 달성"),
        (AMBER,   "지도학습 도메인 갭",
         "합성 노이즈 != 실제 스페클\nPSNR 22.43 (막대 그래프 제외)"),
    ]
    for i, (col, title, body) in enumerate(findings):
        add_rect(s, Inches(9.8), Inches(2.2 + i * 1.52), Inches(0.12), Inches(1.0), col)
        add_text(s, title, Inches(10.02), Inches(2.2 + i * 1.52), Inches(2.85), Inches(0.45),
                 size=11, bold=True, color=ON_DARK, font="나눔고딕")
        add_text(s, body, Inches(10.02), Inches(2.68 + i * 1.52), Inches(2.85), Inches(0.65),
                 size=10, color=ON_DARK_S, font="나눔고딕")

    add_divider(s, Inches(9.8), Inches(6.78), Inches(3.0), CORAL)
    add_text(s, "=> 사전 학습 SR 모델이 단일 이미지에서\n   SRAD 초과 달성. 전체 평가 필요.",
             Inches(9.8), Inches(6.88), Inches(3.05), Inches(0.55),
             size=10, bold=True, color=ON_DARK, font="나눔고딕")


def slide_next(prs: Presentation) -> None:
    """슬라이드 13: 분석 및 향후 계획"""
    s = blank_slide(prs)
    set_bg(s, DARK)

    add_text(s, "분석 및 향후 계획", Inches(0.6), Inches(0.25), Inches(11), Inches(0.8),
             size=30, bold=False, color=ON_DARK, font="나눔고딕")
    add_divider(s, Inches(0.6), Inches(1.1), Inches(12.0), CORAL)

    add_rect(s, Inches(0.4), Inches(1.3), Inches(5.8), Inches(5.8), DARK_ELV)
    add_text(s, "현재 한계 및 확인된 문제", Inches(0.6), Inches(1.5), Inches(5.4), Inches(0.48),
             size=15, bold=True, color=AMBER, font="나눔고딕")
    limits = [
        ("합성 노이즈 도메인 갭",
         "Gamma+Gaussian 합성 노이즈 != 실제 OCT 스페클\n지도학습 SSIM 0.333으로 크게 하락"),
        ("N2N 학습 데이터 부족",
         "39세트(156프레임)는 절대적으로 부족\nloss 0.105 -> 0.099로 수렴 한계"),
        ("SR 단일 이미지 테스트",
         "Real-ESRGAN 1장 테스트만 완료\n18쌍 전체 평가 및 파이프라인 연결 필요"),
        ("치주 도메인 미검증",
         "현재까지 모든 평가가 망막 SBSDI D1 기준\n치주 OCT 전용 평가 프로토콜 필요"),
    ]
    for i, (title, desc) in enumerate(limits):
        add_rect(s, Inches(0.6), Inches(2.1 + i * 1.12), Inches(0.08), Inches(0.88), AMBER)
        add_text(s, title, Inches(0.8), Inches(2.1 + i * 1.12), Inches(5.0), Inches(0.38),
                 size=12, bold=True, color=ON_DARK, font="나눔고딕")
        add_text(s, desc, Inches(0.8), Inches(2.52 + i * 1.12), Inches(5.0), Inches(0.6),
                 size=11, color=ON_DARK_S, font="나눔고딕")

    add_rect(s, Inches(6.6), Inches(1.3), Inches(6.3), Inches(5.8), DARK_ELV)
    add_text(s, "향후 계획", Inches(6.8), Inches(1.5), Inches(5.9), Inches(0.48),
             size=15, bold=True, color=TEAL, font="나눔고딕")

    add_rect(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(0.72), RGBColor(0x1a, 0x19, 0x17))
    add_text(s, "y = D(H(x) * Ns + Na)\nx: 고해상도 clean  D: 다운샘플  H: 측방향 블러",
             Inches(6.9), Inches(2.16), Inches(5.7), Inches(0.62),
             size=11, color=ON_DARK, font="Courier New")

    plans = [
        (CORAL,   "Real-ESRGAN 18쌍 전체 평가",
         "단일 이미지 -> 전체 벤치마크 비교"),
        (TEAL,    "N2N + SR 파이프라인 연결",
         "기존 N2N 디노이징 후 SR 적용 비교"),
        (AMBER,   "PSCAT / PnP-DM 구현",
         "End-to-End Joint 모델\nPSCAT PSNR 31.48 dB @ PKU37 x4"),
        (SUCCESS, "치주 OCT 데이터 확보",
         "교수님과 촬영 프로토콜 논의\n전이 성능 정량 평가"),
    ]
    for i, (col, title, desc) in enumerate(plans):
        add_rect(s, Inches(6.8), Inches(2.98 + i * 1.0), Inches(0.1), Inches(0.78), col)
        add_text(s, title, Inches(7.02), Inches(2.98 + i * 1.0), Inches(5.5), Inches(0.38),
                 size=12, bold=True, color=ON_DARK, font="나눔고딕")
        add_text(s, desc, Inches(7.02), Inches(3.4 + i * 1.0), Inches(5.5), Inches(0.52),
                 size=11, color=ON_DARK_S, font="나눔고딕")

    add_divider(s, Inches(0.6), Inches(7.05), Inches(12.1), CORAL)
    add_text(s, "목표: 스페클 노이즈 제거 + SR 통합 AI 파이프라인 구현 => 치주 OCT 진단 품질 향상",
             Inches(0.6), Inches(7.1), Inches(12.0), Inches(0.35),
             size=13, bold=True, color=ON_DARK, align=PP_ALIGN.CENTER, font="나눔고딕")


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------

def main() -> None:
    prs = new_prs()

    print("슬라이드 생성 중...")
    slide_title(prs);      print("  1/13 표지")
    slide_background(prs); print("  2/13 연구 배경")
    slide_oct_noise(prs);  print("  3/13 OCT와 스페클 노이즈")
    slide_dataset(prs);    print("  4/13 데이터셋")
    slide_metrics(prs);    print("  5/13 평가 지표")
    slide_pipeline(prs);   print("  6/13 파이프라인 개요")
    slide_baseline(prs);   print("  7/13 전통 방법 베이스라인")
    slide_synthetic(prs);  print("  8/13 합성 노이즈 생성")
    slide_n2n(prs);        print("  9/13 자가지도 학습 (N2N)")
    slide_supervised(prs); print(" 10/13 지도학습 (도메인 갭)")
    slide_sr_test(prs);    print(" 11/13 Real-ESRGAN SR 테스트")
    slide_comparison(prs); print(" 12/13 종합 성능 비교")
    slide_next(prs);       print(" 13/13 향후 계획")

    out = RESULTS_DIR / "presentation.pptx"
    prs.save(str(out))
    print(f"\n저장 완료: {out}")


if __name__ == "__main__":
    main()
